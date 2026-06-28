"""
牛马助手（GrindPal）- FastAPI 后端入口 v0.8

安全增强：
- JWT 鉴权保护所有功能端点
- 用户级频率限制（回退 IP）
- 所有 SQL 参数化查询（database.py）
- API Key 仅通过 Header 传递，不落库、不记录日志
"""

import asyncio
import os
import time
import uuid
import re
import json
from datetime import datetime
from contextlib import asynccontextmanager

# PaddlePaddle 3.x ONEDNN 兼容性：禁用 MKLDNN 避免 crash
os.environ.setdefault("FLAGS_use_mkldnn", "0")
os.environ.setdefault('FFMPEG_ARGS', '-c:a pcm_s16le')

from dotenv import load_dotenv
load_dotenv()

from fastapi import FastAPI, HTTPException, Request, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel, Field

from database import (
    init_db, add_history, get_history, get_history_by_id, delete_history,
    create_conversation, list_conversations, get_conversation,
    get_conversation_title, rename_conversation, delete_conversation,
    add_message, get_messages, delete_last_assistant_message,
    delete_last_message_pair, truncate_messages_after, truncate_messages_from, delete_message_by_id,
    update_message_content,
    add_attachment, get_attachments, search_messages,
)
from version import VERSION
from auth import get_current_user
from routers.auth import router as auth_router
from llm import (
    call_llm,
    call_llm_stream,
    build_summarize_prompt,
    build_email_prompt,
    build_minutes_prompt,
    build_polish_prompt,
    build_report_ese_prompt,
    build_requirements_prompt,
    build_prd_prompt,
    build_prd_demo_prompt,
    build_ppt_outline_prompt,
    build_weekly_research_prompt,
    build_weekly_project_prompt,
    build_weekly_techsurvey_prompt,
    build_weekly_ops_prompt,
    generate_conversation_title,
    get_chat_persona,
    ENV_BASE_URL, LLM_STREAM_TIMEOUT,
    SAFETY_GUARD, CHAT_SAFETY_GUARD, CHAT_SYSTEM, is_mock_mode,
)
from mock_responses import mock_response
from openai import AsyncOpenAI

from logger import init_logging, get_logger
from backend_i18n import t as _t

init_logging(os.getenv("LOG_LEVEL", "INFO"))
logger = get_logger("api")
_whisper_model = None  # openai-whisper
_whisper_lock = None    # openai-whisper 初始化锁
_wcp_model = None       # whisper.cpp
_wcp_current_model_name = None  # 当前加载的模型名（tiny/base/small/medium），用于热切换
_wcp_lock = None        # whisper.cpp 并发保护
_webm_header = None      # Chrome MediaRecorder 第一个 chunk 的 EBML 头
# GPU 加速已禁用
_n_gpu_layers = 0
# _n_gpu_layers = int(os.getenv("WHISPER_N_GPU_LAYERS", "0"))
_last_text = {}          # user_id → 上一次转写文本（用于上下文传递+去重）
_last_text_lock = asyncio.Lock()  # 保护 _last_text 并发读写

# ---- 频率限制（用户级优先，回退 IP） ----
RATE_LIMIT = 100  # 每用户每分钟
RATE_WINDOW = 60
MAX_RATE_ENTRIES = 10000  # 最大内存条目数，防止恶意 IP 切换耗尽内存
_rate_map: dict[str, list[float]] = {}
_rate_lock = asyncio.Lock()
_rate_cleanup_counter = 0


async def _check_rate_limit(user_id: int, ip: str, request: Request = None):
    key = f"u{user_id}" if user_id else f"ip{ip}"
    now = time.time()
    window_start = now - RATE_WINDOW
    async with _rate_lock:
        if key in _rate_map:
            _rate_map[key] = [t for t in _rate_map[key] if t > window_start]
            if len(_rate_map[key]) >= RATE_LIMIT:
                raise HTTPException(status_code=429, detail=_t("rate_limit", request) if request else "请求过于频繁，请稍后再试")
            _rate_map[key].append(now)
        else:
            # 超过上限时做一次全局清理，仍然超限则拒绝新条目
            if len(_rate_map) >= MAX_RATE_ENTRIES:
                expired_all = [k for k, v in _rate_map.items() if not v or all(t <= window_start for t in v)]
                for k in expired_all:
                    del _rate_map[k]
                if len(_rate_map) >= MAX_RATE_ENTRIES:
                    raise HTTPException(status_code=429, detail=_t("rate_limit", request) if request else "请求过于频繁，请稍后再试")
            _rate_map[key] = [now]
        # 每 100 次请求清理一次过期 key
        global _rate_cleanup_counter
        _rate_cleanup_counter += 1
        if _rate_cleanup_counter >= 100:
            _rate_cleanup_counter = 0
            expired = [k for k, v in _rate_map.items() if not v or all(t <= window_start for t in v)]
            for k in expired:
                del _rate_map[k]


# ---- Pydantic 模型 ----

class SummarizeRequest(BaseModel):
    text: str = Field(..., max_length=50000)
    length: str = "medium"
    custom_instruction: str = ""

class EmailRequest(BaseModel):
    recipient: str = Field(..., max_length=200)
    subject_keywords: str = Field(..., max_length=200)
    points: list[str] = Field(default_factory=list)
    tone: str = "formal"
    original_email: str = ""
    custom_instruction: str = ""

class MinutesRequest(BaseModel):
    transcript: str = Field(..., max_length=50000)
    speaker_tags: bool = False
    custom_instruction: str = ""

class PolishRequest(BaseModel):
    draft: str = Field(..., max_length=50000)
    style: str = "business"
    custom_instruction: str = ""

class ReportEseRequest(BaseModel):
    rant: str = Field(..., max_length=50000)
    style: str = "result-oriented"
    custom_instruction: str = ""

class RequirementsRequest(BaseModel):
    text: str = Field(..., max_length=50000)
    style: str = "spec"
    custom_instruction: str = ""

class PrdRequest(BaseModel):
    idea: str = Field(..., max_length=50000)
    style: str = "full"
    with_demo: bool = False
    custom_instruction: str = ""

class PptOutlineRequest(BaseModel):
    topic: str = Field(..., max_length=20000)
    points: str = Field(..., max_length=50000)
    style: str = "outline"
    custom_instruction: str = ""


class WeeklyReportRequest(BaseModel):
    """周报生成请求"""
    report_type: str = Field(..., description="research|project|techsurvey|ops")
    style: str = Field("structured", description="structured|narrative")
    lang: str = Field("zh", description="zh|en")
    raw_notes: str = Field("", description="用户原始笔记/素材，LLM自行拆解组织")
    custom_instruction: str = ""


# ---- FastAPI 应用 ----

@asynccontextmanager
async def lifespan(app: FastAPI):
    # 启动时环境变量校验
    from auth import _PLACEHOLDER_SECRETS as _jwt_placeholders
    _jwt = os.getenv("JWT_SECRET", "")
    if _jwt in _jwt_placeholders:
        logger.warning("JWT_SECRET 未设置或为占位值，auth.py 将自动生成随机密钥", extra={"request_id": "system"})
    from llm import _PLACEHOLDER_KEYS
    _api_key = os.getenv("DEEPSEEK_API_KEY", "")
    if not _api_key or _api_key in _PLACEHOLDER_KEYS:
        logger.warning("DEEPSEEK_API_KEY 未配置，LLM 调用将使用 Mock 模式（返回模拟数据）", extra={"request_id": "system"})
    await init_db()
    logger.info(f"GrindPal v{VERSION} 启动，数据库已初始化", extra={"request_id": "system"})
    # OCR 可选依赖检测
    # OCR 暂不可用（PaddleOCR/PaddlePaddle 版本兼容问题，后续修复）
    logger.info("OCR 暂不可用（依赖兼容问题）", extra={"request_id": "system"})
    yield

app = FastAPI(title="GrindPal API", version=VERSION, lifespan=lifespan)

# CORS 配置（生产环境应限定域名）
_CORS_ORIGINS = [o.strip() for o in os.getenv("CORS_ORIGINS", "http://localhost:3000,http://localhost:8000").split(",")]
app.add_middleware(
    CORSMiddleware,
    allow_origins=_CORS_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type", "X-Api-Key", "X-Model", "X-Custom-Instruction", "X-Deep-Think", "X-Style", "X-Kb-Collection"],
)


@app.middleware("http")
async def log_requests(request: Request, call_next):
    """请求日志中间件。注意：不记录 Header 内容，避免 API Key 泄露。uvicorn access_log 已关闭。"""
    request_id = str(uuid.uuid4())[:8]
    request.state.request_id = request_id
    t0 = time.perf_counter()
    response = await call_next(request)
    elapsed = time.perf_counter() - t0
    logger.info(
        f"{request.method} {request.url.path} -> {response.status_code} ({elapsed:.3f}s)",
        extra={"request_id": request_id},
    )
    return response


# ---- 工具函数 ----

def _success(data, message="success", request=None):
    msg = _t(message, request) if isinstance(message, str) and request else message
    return {"code": 200, "data": data, "message": msg}

def _error(code: int, message: str, request=None):
    msg = _t(message, request) if request else message
    return JSONResponse(status_code=code, content={"code": code, "data": None, "message": msg})

def _safe_error(code: int, key: str, e: Exception, request=None):
    """记录真实错误到日志，返回不泄露细节的通用消息"""
    import traceback
    logger.error(f"{key}: {e}\n{traceback.format_exc()}", extra={"request_id": "-"})
    msg = _t(key, request) if request else key
    return _error(code, msg, request)

MAX_UPLOAD_SIZE = 200 * 1024 * 1024  # 200MB 上传上限
KB_SEARCH_LIMIT = 8  # 知识库检索返回片段数

def _parse_kb_collections(header_value: str) -> list[int] | None:
    """解析 X-Kb-Collection 头：逗号分隔的整数 ID 列表，空则返回 None"""
    if not header_value or not header_value.strip():
        return None
    ids = []
    for part in header_value.split(","):
        part = part.strip()
        if part:
            try:
                ids.append(int(part))
            except ValueError:
                pass
    return ids if ids else None

def _kb_overlap_fallback(kb_chunks: list[dict], response_text: str) -> list[dict]:
    """当 LLM 未显式引用 KB 时，通过文本重叠检测哪些片段被隐式使用。
    使用 2-gram 重叠：片段与回复共享 ≥3 个 2-gram 则认为可能被引用。
    """
    if not kb_chunks or not response_text:
        return []
    # 构建回复的 2-gram 集合
    resp_2grams = set()
    cleaned = response_text
    for i in range(len(cleaned) - 1):
        bigram = cleaned[i:i+2].strip()
        if len(bigram) >= 2:
            resp_2grams.add(bigram)
    if not resp_2grams:
        return []
    result = []
    for i, c in enumerate(kb_chunks):
        content = c.get("content", "")
        if not content:
            continue
        # 计算共享 2-gram 数量
        shared = 0
        seen = set()
        for j in range(len(content) - 1):
            bg = content[j:j+2].strip()
            if len(bg) < 2 or bg in seen:
                continue
            seen.add(bg)
            if bg in resp_2grams:
                shared += 1
        # 至少 3 个共享 2-gram 才认为有重叠
        if shared >= 3:
            result.append({**c, "chunk_index": i, "_fallback": True})
    # 最多返回 5 个
    return result[:5]


def _filter_referenced_chunks(kb_chunks: list[dict], response_text: str) -> list[dict]:
    """从 LLM 回复中提取实际引用的参考编号，只返回被引用的片段。
    支持 LLM 输出的多种引用格式变体：
      【参考 1 · 来源：xxx.docx】
      【参考1·来源：xxx】
      （参考 1：来源 xxx）
      [参考1: 来源 xxx]
      参考 1: xxx
      📎 参考1 · xxx
      参考资料 1：xxx
    如果未引用任何片段，返回空列表。
    """
    if not kb_chunks or not response_text:
        return []

    ref_nums = set()

    # 多种正则模式，覆盖 LLM 常见的引用格式变体
    patterns = [
        # 【参考 1 · 来源：xxx】 / 【参考1·来源】
        r'【参考\s*(\d+)',
        # （参考 1：来源） / （参考1: 来源）
        r'[（(]参考\s*(\d+)',
        # [参考1: 来源] / [参考 1：来源]
        r'\[参考\s*(\d+)',
        # 参考 1: / 参考1：/ 参考 1 ·
        r'(?:^|\n|。|，)\s*参考\s*(\d+)\s*[：:·]',
        # 📎 参考1 · / 📎1·
        r'📎\s*(?:参考\s*)?(\d+)\s*[·：:]',
        # 参考资料 1： / 引用 1:
        r'(?:参考资料|引用|文献)\s*(\d+)\s*[：:·]',
        # ——参考 5 / — 参考 5 / ---参考 5
        r'[—\-]{2,3}\s*参考\s*(\d+)',
    ]
    for pat in patterns:
        for m in re.finditer(pat, response_text):
            ref_nums.add(int(m.group(1)))

    # 也尝试从文末引用块中直接匹配文件名
    # 格式: 【参考 N · 来源：filename.docx】
    cited_filenames = set()
    for m in re.finditer(r'【参考\s*\d+\s*[·\s]*来源[：:]\s*([^】]+)】', response_text):
        cited_filenames.add(m.group(1).strip())

    if not ref_nums and not cited_filenames:
        # LLM 没有显式引用任何 KB 片段 → 不通过文本重叠猜测，直接返回空
        return []

    # 先按编号过滤（1-based → 0-based）
    result_by_idx = {}
    for i, c in enumerate(kb_chunks):
        if (i + 1) in ref_nums:
            result_by_idx[i] = c

    # 再按文件名过滤（补充编号未匹配到的）
    for i, c in enumerate(kb_chunks):
        if i not in result_by_idx:
            fn = c.get("filename", "")
            if any(cf in fn or fn in cf for cf in cited_filenames):
                result_by_idx[i] = c

    # 无法精确匹配任何引用 → 返回空，不猜测
    if not result_by_idx:
        return []

    # 按 chunk_index 排序返回
    return [result_by_idx[i] for i in sorted(result_by_idx.keys())]


def _kb_results_relevant(chunks: list[dict], keyword_count: int) -> bool:
    """检查 KB 检索结果是否有足够的相关性。
    要求最佳片段至少命中 ceil(keyword_count/2) 个关键词，避免常见词误匹配。
    """
    if not chunks:
        return False
    # 过滤掉内容过短（<10字）或 score 为 0 的噪声片段
    usable = [c for c in chunks
              if c.get("score", 0) >= 1 and len(c.get("content", "").strip()) >= 10]
    if not usable:
        return False
    import math
    best_score = max((c.get("score", 0) for c in usable), default=0)
    # 至少需要 ceil(kw/2) 个命中，且最低 2 个（防止单关键词偶然匹配）
    min_required = max(2, math.ceil(keyword_count / 2))
    if best_score < min_required:
        return False
    return True


def _filter_usable_chunks(chunks: list[dict]) -> list[dict]:
    """过滤 KB 检索结果：只保留 score >= 1 且内容 >= 10 字的片段，按 score 降序"""
    usable = [c for c in chunks
              if c.get("score", 0) >= 1 and len(c.get("content", "").strip()) >= 10]
    usable.sort(key=lambda c: c.get("score", 0), reverse=True)
    return usable

def _check_upload_size(request: Request):
    """检查 Content-Length，超限返回 413"""
    cl = request.headers.get("Content-Length")
    if cl:
        try:
            size = int(cl)
        except (ValueError, TypeError):
            return _error(400, _t("bad_request", request))
        if size > MAX_UPLOAD_SIZE:
            return _error(413, _t("file_too_large", request, size=size//1024//1024, max=MAX_UPLOAD_SIZE//1024//1024))
    return None

def _get_headers(request: Request) -> tuple[str, str, str]:
    api_key = request.headers.get("X-Api-Key", "")
    model = request.headers.get("X-Model", "")
    custom = request.headers.get("X-Custom-Instruction", "")
    deep_think = request.headers.get("X-Deep-Think", "0") == "1"
    return api_key, model, custom, deep_think

def _redact_api_key(text: str) -> str:
    """将 API Key 替换为 [REDACTED]，防止日志泄露"""
    if not text:
        return text
    # 匹配 sk- 开头或类似 API Key 模式
    return re.sub(r'(sk-[a-zA-Z0-9_-]{4})[a-zA-Z0-9_-]+', r'\1...[REDACTED]', str(text))

def _get_style(request: Request) -> str:
    """提取文风偏好: natural | standard | formal"""
    return request.headers.get("X-Style", "standard")

KB_KEYWORD_SYSTEM = (
    "你是关键词提取助手。分析用户消息的语义，提取 3-10 个用于全文搜索的关键词或短语。\n"
    "规则：\n"
    "1. 优先提取专有名词、术语、人名、地名、日期\n"
    "2. 补充同义词和相关概念（如用户问「怎么瘦肚子」，提取：瘦肚子、腹部减脂、腰腹运动）\n"
    "3. 忽略无意义的虚词和口语填充词\n"
    "4. 只输出逗号分隔的关键词，一行即可，不要任何解释"
)

async def _extract_kb_keywords(user_message: str, api_key: str, model: str) -> list[str]:
    """通过 LLM 提取知识库搜索关键词，失败时回退到 jieba 分词"""
    # 优先使用 LLM 提取语义关键词（排除占位/无效 Key，直接走 jieba）
    if api_key and not is_mock_mode(api_key):
        try:
            client = AsyncOpenAI(
                api_key=api_key,
                base_url=ENV_BASE_URL,
                timeout=10.0,
            )
            try:
                resp = await client.chat.completions.create(
                    model=model or "deepseek-chat",
                    messages=[
                        {"role": "system", "content": KB_KEYWORD_SYSTEM},
                        {"role": "user", "content": user_message[:2000]},
                    ],
                    max_tokens=120,
                    temperature=0,
                )
                text = resp.choices[0].message.content.strip()
                # 解析逗号/顿号分隔的关键词
                keywords = [k.strip() for k in text.replace('，', ',').replace('、', ',').split(',') if k.strip()]
                if keywords:
                    logger.info(f"LLM keywords: {keywords}")
                    return [k for k in keywords if len(k) >= 1][:20]
            finally:
                await client.close()
        except Exception as e:
            logger.warning(f"LLM keyword extraction failed, falling back to jieba: {e}")

    # 回退：jieba 分词
    try:
        import jieba
        words = list(jieba.cut(user_message))
        return [w for w in words if len(w) >= 2 and not w.isspace()][:20]
    except Exception:
        # 最后兜底：简单空格分割
        return [w for w in user_message.split() if len(w) >= 2][:20]

# 中文停用词/疑问词——当关键词全落入此集合时跳过知识库检索
_KB_SKIP_WORDS = frozenset({
    "你", "我", "他", "她", "它", "你们", "我们", "他们",
    "是", "的", "了", "吗", "呢", "吧", "啊", "呀", "哦", "嗯",
    "什么", "怎么", "怎样", "为什么", "哪里", "哪儿", "谁",
    "这", "那", "这个", "那个", "这些", "那些",
    "可以", "能", "会", "要", "想", "应该", "必须",
    "如果", "因为", "所以", "但是", "而且", "或者", "虽然", "然后", "还是",
    "一个", "一下", "一些", "有点", "没有", "什么", "怎么",
    "请问", "帮忙", "告诉", "知道", "觉得", "认为",
    "你好", "谢谢", "再见", "请问",
    "介绍", "如何", "怎么用", "能不能", "可不可以", "做什么",
})

def _need_kb_search(user_message: str) -> bool:
    """快速预判：用户消息是否可能涉及知识库内容。
    在昂贵的 LLM 关键词提取之前运行，用本地规则过滤明显不需要 KB 的消息，
    避免闲聊/问候/确认/管理操作等触发无意义的 KB 检索。
    """
    msg = user_message.strip()
    if not msg:
        return False
    # 极短消息不触发 KB（< 5 字基本是应答或问候）
    if len(msg) < 5:
        return False
    # 纯闲聊/问候/感谢/确认 模式
    casual_patterns = [
        r'^(你好|您好|hi|hello|嗨|嘿|谢谢|多谢|感谢|辛苦了|再见|拜拜|好的|OK|ok|嗯|哦|对|是的|没错|知道了|明白了|收到)[\s!！。.~～,，]*$',
        r'^(早上好|下午好|晚上好|早安|晚安|午安)[\s!！。.~～]*$',
        r'^(可以|能|会|要|行|好|是|对|有)[\s!！。.~～]*$',
        r'^(不|没|不是|没有|算了|不用|不需要)[\s!！。.~～]*$',
        r'^[\s!！。.，,、？?！!~～….\-—]+$',  # 纯标点/空白
        r'^(你是谁|你能做什么|你有什么功能|你叫什么|介绍一下自己|自我介绍)[？?！!。.]*$',
        r'^(继续|接着|然后呢|还有呢|再来|下一个|换一个|再说|多说|接着说)[\s!！。.~～]*$',
        r'^(帮我|请帮我|能不能帮我|可以帮我|麻烦帮我)[\s]*$',  # 只有"帮我"没有具体内容
    ]
    for pat in casual_patterns:
        if re.match(pat, msg):
            return False
    # 纯管理操作：设置、切换、查看配置等
    admin_patterns = [
        r'^(切换|设置|修改|更改|调整|打开|关闭|启用|禁用|查看|显示|列出).{0,8}(模型|主题|模式|风格|语言|人设|知识库|模板|历史|对话|会话)',
        r'^(新建|删除|清空|清除|清理|导出|导入|备份|恢复)',
        r'^(设置|配置|选项|偏好|参数|首选项)',
        r'^(翻到|回到|跳到|滚动|上一页|下一页|前面|后面)',
    ]
    for pat in admin_patterns:
        if re.match(pat, msg):
            return False
    # 图片/文件查询：如果用户消息主要围绕某个文件名或图片链接，不需要 KB 检索
    file_query_patterns = [
        r'\b[\w.-]+\.(jpg|jpeg|png|gif|webp|bmp|svg|pdf|docx|xlsx|pptx|zip|rar)\b',
        r'(这张|那个|这个|看看|帮我看看|分析一下|识别).{0,6}(图|图片|照片|截图|文件|文档)',
        r'(图|图片|照片|画|文件|截图).{0,4}(是什么|里面|内容|写了|识别)',
        r'(pixiv|插画|漫画|同人|壁纸|头像|QQ\d*)',
    ]
    for pat in file_query_patterns:
        if re.search(pat, msg):
            return False
    return True


def _should_search_kb(keywords: list[str]) -> bool:
    """判断关键词是否值得触发知识库检索"""
    if not keywords or len(keywords) < 2:
        return False
    # 如果所有关键词都是停用词/疑问词，跳过
    if all(kw in _KB_SKIP_WORDS for kw in keywords):
        return False
    # 如果超过一半关键词是停用词，说明用户问题太泛化（如闲聊、自我介绍询问），跳过
    stop_count = sum(1 for kw in keywords if kw in _KB_SKIP_WORDS)
    if stop_count > len(keywords) / 2:
        return False
    return True

def _style_temperature(style: str) -> float:
    """文风 → temperature 映射"""
    return {"natural": 0.85, "standard": 0.7, "formal": 0.5}.get(style, 0.7)

def _style_instruction(style: str) -> str:
    """文风 → 追加系统指令"""
    if style == "natural":
        return "\n\n== 文风模式：自然口语 ==\n用日常对话语气写作，像在给信任的同事发消息。允许口语词、短句、偶尔的情绪表达。不要写得很'正式'，但要保持清晰。"
    elif style == "formal":
        return "\n\n== 文风模式：正式公文体 ==\n用正式但不僵硬的语言。禁止模板化套话（如'综上所述''在……背景下'），但整体保持职业和严谨。"
    return ""  # standard = default, no extra instruction

# ---- 关键词前置过滤 ----
KEYWORD_BLOCK = [
    # 越狱 / 提示词泄露
    "忽略指令", "ignore instructions", "忽略规则", "ignore all",
    "系统提示词", "system prompt", "show me your system",
    "output your prompt", "开发者模式", "DAN模式", "dan模式",
    "越狱", "jailbreak", "STAN模式", "角色扮演",
    # 常见攻击 payload
    "[INST]", "[/INST]", "system:", "System:",
    "假装", "虚构一个", "从现在开始你是",
    "you are now", "pretend to be",
    # 其他注入尝试
    "忽略上述", "disregard above", "ignore above",
    "do anything now", "进入开发者模式",
]
# 零宽字符检测（用于防止关键词绕过）
_ZERO_WIDTH_CHARS = re.compile(r'[\u200B\u200C\u200D\uFEFF\u200E\u200F\u202A-\u202E\u2060-\u2064]')

def _check_keywords(text: str, request: Request):
    lower = text.lower()
    for kw in KEYWORD_BLOCK:
        if kw.lower() in lower:
            rid = getattr(request.state, "request_id", "-")
            logger.warning(f"关键词拦截 rid={rid} keyword={kw}", extra={"request_id": rid})
            raise HTTPException(status_code=400, detail=_t("blocked_content", request))
    # 独立 DAN（Do Anything Now）单词级检测
    if re.search(r'\bdan\b', lower):
        rid = getattr(request.state, "request_id", "-")
        logger.warning(f"关键词拦截(DAN单词) rid={rid}", extra={"request_id": rid})
        raise HTTPException(status_code=400, detail=_t("blocked_content", request))
    # 零宽字符检测
    if _ZERO_WIDTH_CHARS.search(text):
        rid = getattr(request.state, "request_id", "-")
        logger.warning(f"零宽字符检测拦截 rid={rid}", extra={"request_id": rid})
        raise HTTPException(status_code=400, detail=_t("blocked_content", request))

# ---- 敏感信息脱敏 ----
def _mask_sensitive(text: str) -> str:
    text = re.sub(r'1[3-9]\d{9}', '[手机号已隐藏]', text)
    text = re.sub(r'\d{17}[\dXx]', '[身份证号已隐藏]', text)
    text = re.sub(r'\b\d{16,19}\b', '[银行卡号已隐藏]', text)
    text = re.sub(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', '[邮箱已隐藏]', text)
    return text

# ---- Router 初始化 ----
from routers.auth import init as auth_init
auth_init(_check_rate_limit, _success, _error)
app.include_router(auth_router)

async def _call_and_save(
    type_: str, input_text: str, system_prompt: str, user_prompt: str,
    request: Request, user_id: int, api_key: str = "", model: str = "",
    temperature: float = 0.7, max_tokens: int = 20000, custom_instruction: str = "",
) -> dict:
    if custom_instruction:
        system_prompt = system_prompt + "\n\n【用户自定义指令】\n" + custom_instruction
    # 知识库检索注入
    kb_chunks = []
    kb_col = request.headers.get("X-Kb-Collection", "")
    kb_cids = _parse_kb_collections(kb_col) if kb_col else None
    if kb_cids and _need_kb_search(user_prompt):
        try:
            keywords = await _extract_kb_keywords(user_prompt, api_key, model)
            if keywords and _should_search_kb(keywords):
                from database import kb_search
                query = " OR ".join(keywords)
                chunks = await kb_search(user_id, query, limit=KB_SEARCH_LIMIT, collection_ids=kb_cids)
                chunks = _filter_usable_chunks(chunks)
                if chunks and _kb_results_relevant(chunks, len(keywords)):
                    kb_chunks = [{"filename": c["filename"], "content": c["content"], "chunk_index": i}
                                 for i, c in enumerate(chunks)]
                    refs = "\n".join(f"【参考 {i+1} · 来源：{c['filename']}】\n{c['content'][:800]}"
                                    for i, c in enumerate(chunks))
                    system_prompt = system_prompt + "\n\n【知识库参考资料】\n" + refs + \
                        "\n\n以上资料根据关键词自动检索，仅供参考。请只关注用户本次发送的最新消息——只有当资料与本条消息直接相关时才引用。不要因为对话历史里聊过相关话题就引用资料，历史话题与当前问题可能已无关。引用格式：【参考 编号 · 来源：文件名】。如果不相关，直接忽略。"
        except Exception:
            logger.warning(f"KB search skipped in _call_and_save")
    # 关键词前置过滤
    _check_keywords(user_prompt, request)
    # 敏感信息脱敏
    user_prompt = _mask_sensitive(user_prompt)
    # 输入语言检测
    from backend_i18n import detect_input_lang
    input_lang = detect_input_lang(user_prompt)
    if input_lang == "zh":
        system_prompt = "== 重要指令 ==\n用户输入为中文，你的全部输出必须使用中文。\n\n" + system_prompt
    else:
        system_prompt = "== IMPORTANT ==\nThe user input is in English. You MUST respond entirely in English. Do NOT output any Chinese.\n\n" + system_prompt
    try:
        content, usage = await call_llm(system_prompt, user_prompt, api_key=api_key, model=model, temperature=temperature, max_tokens=max_tokens)
        content = _mask_sensitive(content)  # 输出脱敏：防止 LLM 复述输入中的敏感信息
    except Exception as e:
        rid = getattr(request.state, "request_id", "-")
        logger.error(f"AI调用失败 type={type_} user={user_id} error={str(e)}", extra={"request_id": rid})
        raise HTTPException(status_code=502, detail=_t("ai_error", request))
    record_id = await add_history(user_id, type_, input_text, content, tokens=usage.get("total_tokens", 0))
    # 只返回 LLM 实际引用的 KB 片段
    _kb_chunks_out = _filter_referenced_chunks(kb_chunks, content)
    return {"result": content, "usage": usage, "record_id": record_id, "kb_chunks": _kb_chunks_out}


@app.get("/api/v1/templates")
async def list_templates(user: dict = Depends(get_current_user)):
    from database import get_templates
    templates = await get_templates(user["id"])
    return _success({"templates": templates})


@app.post("/api/v1/templates")
async def create_template_api(request: Request, user: dict = Depends(get_current_user)):
    from database import create_template as db_create
    import json as _json
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    name = str(body.get("name", "")).strip()
    if not name:
        return _error(400, _t("template_name_required", request))
    tid = await db_create(
        user["id"], name,
        _json.dumps(body.get("modules", []), ensure_ascii=False),
        str(body.get("system_prompt", "")),
        str(body.get("output_style", "paragraph")),
        int(body.get("is_default", 0)),
    )
    if tid is None:
        return _error(500, _t("template_create_failed", request))
    return _success({"id": tid}, "template_created", request)


@app.put("/api/v1/templates/{template_id}")
async def update_template_api(template_id: int, request: Request, user: dict = Depends(get_current_user)):
    from database import update_template as db_update
    import json as _json
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    kwargs = {}
    for k in ("name", "system_prompt", "output_style", "is_default"):
        if k in body:
            kwargs[k] = body[k]
    if "modules" in body:
        kwargs["modules"] = _json.dumps(body["modules"], ensure_ascii=False)
    ok = await db_update(template_id, user["id"], **kwargs)
    if not ok:
        return _error(404, _t("template_not_found_or_denied", request))
    return _success(None, "template_updated", request)


@app.delete("/api/v1/templates/{template_id}")
async def delete_template_api(template_id: int, request: Request, user: dict = Depends(get_current_user)):
    from database import delete_template as db_delete
    ok = await db_delete(template_id, user["id"])
    if not ok:
        return _error(404, _t("template_delete_denied", request))
    return _success(None, "template_deleted", request)


@app.post("/api/v1/templates/extract")
async def extract_template(request: Request, user: dict = Depends(get_current_user)):
    """一键提炼风格：从历史输出中反推 system_prompt"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    sample = str(body.get("sample_output", "")).strip()
    if not sample:
        return _error(400, _t("template_sample_required", request))
    api_key, model, _, _ = _get_headers(request)
    sys_p = (
        "# Role: 风格提炼师\n"
        "根据提供的输出样例，反推生成一个简洁的系统指令。\n"
        "## Rules\n"
        "1. 指令以「始终」「确保」「避免」等动词开头\n"
        "2. 指令列表不超过 200 字\n"
        "3. 每条指令指向一个具体的行为约束（不写抽象的原则）\n"
        "4. 只输出指令列表，不加任何解释或标题"
    )
    try:
        content, _ = await call_llm(sys_p, f"输出样例：\n\n{sample}", api_key=api_key, model=model, temperature=0.3, max_tokens=300)
        return _success({"system_prompt": (content or "").strip()})
    except Exception as e:
        return _safe_error(502, _t("extract_template_failed", request), e, request)


# ============================================================
#  待办提取
# ============================================================

@app.post("/api/v1/extract-todos")
async def extract_todos(request: Request, user: dict = Depends(get_current_user)):
    """从文本中提取待办事项（仅提取不保存，由前端审核后保存）"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    text = str(body.get("text", "")).strip()
    if not text:
        return _error(400, _t("template_text_required", request))
    api_key, model, custom, _ = _get_headers(request)
    from datetime import datetime as _dt, timedelta as _td
    now = _dt.now()
    now_str = now.strftime("%Y年%m月%d日 %H:%M")
    tomorrow_str = (now + _td(days=1)).strftime("%Y-%m-%d")
    sys_p = (
        "# Role: 待办提取器\n"
        "## Profile\n"
        "你从会议记录/聊天记录中精准提取待办事项。你判断一条内容是不是「待办」的唯一标准：这行字描述的是一个「将来要做的事」而不是「已经发生的事」。\n"
        "## Workflow\n"
        "1. 逐句扫描：区分「陈述事实」（已完成/背景）和「待办事项」（将来要做）\n"
        "2. 对每条待办：确定任务内容、截止时间、负责人\n"
        "3. 输出JSON数组\n"
        "## Rules\n"
        "1. deadline格式：YYYY-MM-DDTHH:MM（如2026-07-15T14:30），只写了日期没写时间的默认18:00，完全没有时间信息填「待定」\n"
        "2. 相对时间根据当前时间推算：「明天」→+1天、「后天」→+2天、「下周X」→下周对应日期、「今晚」→20:00、「明早」→09:00\n"
        "3. assignee只在原文明确指派某人时才填，否则空字符串「」\n"
        "4. 没有待办返回空数组 []\n"
        "5. 只输出JSON数组，不加任何Markdown或解释文字\n"
        f"当前时间：{now_str}\n"
        f"示例——如果当前是 {now_str}，原文说「明天下午3点开会」，deadline 应为 \"{tomorrow_str}T15:00\"。\n\n"
        "每条待办格式：\n"
        '{"task": "任务内容", "deadline": "截止时间", "assignee": "负责人"}\n\n'
    )
    if custom:
        sys_p = sys_p + "\n\n【用户自定义指令】\n" + custom
    try:
        content, _ = await call_llm(sys_p, text, api_key=api_key, model=model, temperature=0.1, max_tokens=1024)
        import json as _json_todo
        raw = (content or "").strip()
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        try:
            todos = _json_todo.loads(raw)
        except Exception:
            return _error(502, _t("ai_format_error_mock", request))
        if isinstance(todos, dict):
            todos = [todos]
        if not isinstance(todos, list):
            return _error(502, _t("ai_format_error", request))
        # 标准化字段，过滤完全空白的 task
        import re as _re
        cleaned = []
        for t in todos[:20]:
            task = str(t.get("task", "")).strip()
            if not task or task.lower() in ("null", "none", "无"):
                continue
            deadline = str(t.get("deadline", "")).strip()
            # 标准化 deadline：统一为 YYYY-MM-DDTHH:MM
            if not deadline or deadline.lower() in ("null", "none", "无"):
                deadline = "待定"
            else:
                m = _re.match(r'^(\d{4}-\d{2}-\d{2})[ T](\d{2}:\d{2})(:\d{2})?$', deadline)
                if m:
                    deadline = f"{m.group(1)}T{m.group(2)}"
                else:
                    deadline = "待定"
            assignee = str(t.get("assignee", "")).strip()
            if assignee.lower() in ("null", "none", "无"):
                assignee = ""
            cleaned.append({"task": task, "deadline": deadline, "assignee": assignee})
        # 不自动保存，返回给前端审核
        return _success({"todos": cleaned})
    except Exception as e:
        return _safe_error(502, _t("extract_todos_failed", request), e, request)


@app.get("/api/v1/todos")
async def list_todos(user: dict = Depends(get_current_user)):
    from database import get_todos
    todos = await get_todos(user["id"])
    return _success({"todos": todos})


@app.put("/api/v1/todos/{todo_id}")
async def update_todo(todo_id: int, request: Request, user: dict = Depends(get_current_user)):
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    from database import update_todo, update_todo_status
    # 如果有 task 字段 → 编辑模式
    if "task" in body:
        ok = await update_todo(todo_id, user["id"],
            task=body.get("task", ""),
            assignee=body.get("assignee", ""),
            deadline=body.get("deadline", ""),
        )
    else:
        status = body.get("status", "")
        if status not in ("pending", "done"):
            return _error(400, _t("todo_status_invalid", request))
        ok = await update_todo_status(todo_id, user["id"], status)
    if not ok:
        return _error(404, _t("todo_not_found", request))
    return _success(None, "template_updated", request)


@app.post("/api/v1/todos")
async def create_todo(request: Request, user: dict = Depends(get_current_user)):
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    task = str(body.get("task", "")).strip()
    if not task:
        return _error(400, _t("todo_content_required", request))
    from database import add_todo
    tid = await add_todo(user["id"],
        str(body.get("assignee", "")),
        task,
        str(body.get("deadline", "")),
        body.get("source_record_id"),
    )
    return _success({"id": tid}, "todo_created", request)


@app.post("/api/v1/parse-todo")
async def parse_todo(request: Request, user: dict = Depends(get_current_user)):
    """智能解析一句话为待办（提取任务、截止时间、负责人）"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    raw = str(body.get("text", "")).strip()
    if not raw:
        return _error(400, _t("todo_input_required", request))
    api_key, model, _, _ = _get_headers(request)
    from datetime import datetime as _dt, timedelta as _td
    now = _dt.now()
    now_str = now.strftime("%Y年%m月%d日 %H:%M")
    sys_p = (
        "# Role: 待办解析器\n"
        "从用户的一句话中提取待办信息，返回JSON：{\"task\":\"任务内容\",\"deadline\":\"YYYY-MM-DDTHH:MM\",\"assignee\":\"负责人\"}\n"
        f"当前时间：{now_str}\n\n"
        "## Rules\n"
        "1. task：去掉时间词和负责人后剩余的核心动作，用口语祈使句\n"
        "2. deadline：精确到分钟（如2026-07-15T14:30），相对时间按当前推算，只写日期默认18:00，无时间信息填\"待定\"\n"
        "3. assignee：只有明确指派时才填人名，否则\"\"\n"
        "4. 只输出JSON对象，不包裹Markdown、不注释"
    )
    try:
        content, _ = await call_llm(sys_p, raw, api_key=api_key, model=model, temperature=0.1, max_tokens=256)
        import json as _json
        raw_out = (content or "").strip()
        if raw_out.startswith("```"): raw_out = raw_out.split("```")[1].replace("json","").strip()
        parsed = _json.loads(raw_out)
        if not isinstance(parsed, dict): parsed = {}
        return _success({
            "task": str(parsed.get("task", raw)),
            "deadline": str(parsed.get("deadline", "")),
            "assignee": str(parsed.get("assignee", "")),
        })
    except Exception as e:
        return _safe_error(502, _t("parse_todo_failed", request), e, request)


@app.delete("/api/v1/todos/{todo_id}")
async def delete_todo(todo_id: int, request: Request, user: dict = Depends(get_current_user)):
    from database import delete_todo
    ok = await delete_todo(todo_id, user["id"])
    if not ok:
        return _error(404, _t("todo_not_found", request))
    return _success(None, "todo_deleted", request)


# ============================================================
#  多轮对话追问
# ============================================================

@app.post("/api/v1/continue")
async def continue_conversation(request: Request, user: dict = Depends(get_current_user)):
    """基于历史记录继续对话，最多回溯 3 轮"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    record_id = body.get("record_id")
    instruction = str(body.get("instruction", "")).strip()
    if not record_id or not instruction:
        return _error(400, _t("record_need_params", request))
    from database import get_history_chain
    chain = await get_history_chain(int(record_id), user["id"], max_depth=3)
    if not chain:
        return _error(404, _t("record_not_found", request))
    # 组装上下文
    messages = []
    for r in chain:
        messages.append({"role": "user", "content": r["input_text"]})
        messages.append({"role": "assistant", "content": r["result_text"]})
    messages.append({"role": "user", "content": instruction})
    api_key, model, custom, _ = _get_headers(request)
    system_msg = (
        "# Role: 智能助手\n"
        "## Profile\n"
        "基于对话历史理解用户追问意图，给出连贯的回答。如果是修正前次输出，明确指出修改了什么。\n"
        "## Rules\n"
        "1. 先回顾历史对话确定上下文\n"
        "2. 如果是追问细节→展开上次回复的具体部分\n"
        "3. 如果是修正错误→承认并给出修正版本\n"
        "4. 如果是新问题→基于历史上下文回答"
    )
    if custom:
        system_msg = system_msg + "\n\n【用户自定义指令】\n" + custom
    try:
        content, usage = await call_llm(system_msg,
                                        "", api_key=api_key, model=model,
                                        temperature=0.7, max_tokens=20000,
                                        messages=messages)
        await add_history(user["id"], chain[0]["type"],
                         f"追问: {instruction}", content,
                         parent_id=int(record_id),
                         tokens=usage.get("total_tokens", 0))
        return _success({"result": content, "usage": usage})
    except Exception as e:
        return _safe_error(502, _t("continue_failed", request), e, request)


# ============================================================
#  知识库
# ============================================================

@app.get("/api/v1/kb/collections")
async def kb_list(user: dict = Depends(get_current_user)):
    from database import kb_list_collections
    cols = await kb_list_collections(user["id"])
    return _success({"collections": cols})


@app.post("/api/v1/kb/collections")
async def kb_create_col(request: Request, user: dict = Depends(get_current_user)):
    from database import kb_create_collection
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    name = str(body.get("name", "")).strip()
    if not name:
        return _error(400, _t("kb_name_required", request))
    cid = await kb_create_collection(user["id"], name)
    if cid is None:
        return _error(500, _t("kb_create_failed", request))
    return _success({"id": cid}, "kb_created", request)


@app.delete("/api/v1/kb/collections/{col_id}")
async def kb_delete_col(col_id: int, request: Request = None, user: dict = Depends(get_current_user)):
    from database import kb_delete_collection
    ok = await kb_delete_collection(col_id, user["id"])
    if not ok:
        return _error(404, _t("kb_not_found", request))
    return _success(None, "kb_deleted", request)



@app.post("/api/v1/kb/upload")
async def kb_upload(request: Request, user: dict = Depends(get_current_user)):
    """上传文档到知识库，支持批量上传（最大 10MB/文件，PDF 最多 500 页）"""
    if size_err := _check_upload_size(request): return size_err
    from database import kb_add_document, _chunk_text
    import asyncio as _asyncio
    form = await request.form()
    files = form.getlist("file")
    col_id = form.get("collection_id")
    if not files or not col_id:
        return _error(400, _t("kb_file_missing", request))
    results = []
    for file in files:
        filename = file.filename or "unknown"
        raw = await file.read()
        MAX_SIZE = 200 * 1024 * 1024
        if len(raw) > MAX_SIZE:
            results.append({"filename": filename, "error": _t("kb_file_too_large", request, size=len(raw)//1024//1024)})
            continue
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        text = ""
        try:
            if ext == 'pdf':
                import fitz
                doc = fitz.open(stream=raw, filetype='pdf')
                if doc.page_count > 500:
                    doc.close()
                    results.append({"filename": filename, "error": _t("kb_pdf_too_many", request, pages=doc.page_count)})
                    continue
                doc.close()
            text = await _asyncio.to_thread(_extract_doc_text, raw, ext, filename)
            if text.startswith('[') and text.endswith(']'):
                results.append({"filename": filename, "error": _t("kb_file_unsupported", request, name=filename)})
                continue
        except Exception as e:
            results.append({"filename": filename, "error": _t("file_parse_failed", request)})
            continue
        if not text.strip():
            results.append({"filename": filename, "error": _t("kb_file_no_text", request)})
            continue
        MAX_TEXT = 200000
        if len(text) > MAX_TEXT:
            text = text[:MAX_TEXT]
        chunks = _chunk_text(text)
        doc_id = await kb_add_document(int(col_id), filename,
                                        filename.rsplit('.', 1)[-1] if '.' in filename else '',
                                        len(raw), chunks, user["id"])
        results.append({"doc_id": doc_id, "filename": filename, "chunks": len(chunks)})
    return _success({"results": results}, "doc_uploaded", request)


@app.post("/api/v1/kb/upload-text")
async def kb_upload_text(request: Request, user: dict = Depends(get_current_user)):
    """粘贴文本入库"""
    from database import kb_add_document
    import json as _json
    form = await request.form()
    collection_id = int(form.get("collection_id", 0))
    filename = form.get("filename", "粘贴文本.txt")
    file_type = form.get("file_type", "text/plain")
    chunks_str = form.get("chunks", "[]")
    try: chunks = _json.loads(chunks_str)
    except (json.JSONDecodeError, TypeError, ValueError): chunks = [form.get("chunks", "")]
    if not collection_id: return _error(400, "缺少 collection_id")
    if not chunks: return _error(400, _t("text_empty", request))
    doc_id = await kb_add_document(collection_id, filename, file_type, len(str(chunks)), chunks, user["id"])
    if not doc_id: return _error(500, _t("text_save_failed", request))
    return _success({"doc_id": doc_id, "chunks": len(chunks)}, "text_saved", request)


@app.get("/api/v1/kb/documents/{col_id}")
async def kb_collection_docs(col_id: int, user: dict = Depends(get_current_user)):
    """列出知识库中的文档"""
    from database import kb_list_docs
    docs = await kb_list_docs(col_id, user["id"])
    return _success({"documents": docs})


@app.delete("/api/v1/kb/documents/{doc_id}")
async def kb_delete_doc(doc_id: int, request: Request = None, user: dict = Depends(get_current_user)):
    """删除知识库中的单个文档"""
    from database import kb_delete_document
    ok = await kb_delete_document(doc_id, user["id"])
    return _success(None, "kb_deleted" if ok else "doc_not_found", request) if ok else _error(404, _t("doc_not_found", request))


@app.get("/api/v1/kb/search")
async def kb_search_endpoint(q: str = "", request: Request = None, user: dict = Depends(get_current_user)):
    """知识库全文检索"""
    if not q.strip():
        return _error(400, _t("search_query_required", request))
    try:
        import jieba
        words = list(jieba.cut(q))
        keywords = [w for w in words if len(w) >= 2 and not w.isspace()][:20]
        if not keywords:
            return _success({"chunks": []})
        from database import kb_search
        # 清洗 FTS5 特殊字符
        safe_keywords = [w.replace('"', '').replace("'", "").replace("*", "") for w in keywords]
        safe_keywords = [w for w in safe_keywords if w]  # 移除空词
        if not safe_keywords:
            return _success({"chunks": []})
        query = " OR ".join(safe_keywords)
        chunks = await kb_search(user["id"], query, limit=KB_SEARCH_LIMIT)
        return _success({"chunks": chunks, "keywords": keywords})
    except ImportError:
        # jieba 未安装时回退到按空格/标点分词
        import re as _re
        words = _re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z0-9]+', q)
        keywords = [w for w in words if len(w) >= 2][:20]
        if not keywords:
            return _success({"chunks": []})
        from database import kb_search
        safe_keywords = [w.replace('"', '').replace("'", "").replace("*", "") for w in keywords]
        safe_keywords = [w for w in safe_keywords if w]
        if not safe_keywords:
            return _success({"chunks": []})
        query = " OR ".join(safe_keywords)
        chunks = await kb_search(user["id"], query, limit=KB_SEARCH_LIMIT)
        return _success({"chunks": chunks, "keywords": keywords})
    except Exception as e:
        return _safe_error(500, _t("search_failed", request), e, request)


# ============================================================
#  功能端点（需 JWT 鉴权）
# ============================================================


async def _stream(
    type_: str, input_text: str, system_prompt: str, user_prompt: str,
    request: Request, user_id: int, api_key: str, model: str,
    temperature: float = 0.7, max_tokens: int = 20000, custom_instruction: str = "",
):
    """SSE 流式响应生成器"""
    if custom_instruction:
        system_prompt = system_prompt + "\n\n【用户自定义指令】\n" + custom_instruction
    # 知识库检索注入
    kb_chunks = []
    kb_col = request.headers.get("X-Kb-Collection", "")
    kb_cids = _parse_kb_collections(kb_col) if kb_col else None
    if kb_cids and _need_kb_search(user_prompt):
        try:
            keywords = await _extract_kb_keywords(user_prompt, api_key, model)
            if keywords and _should_search_kb(keywords):
                from database import kb_search
                query = " OR ".join(keywords)
                chunks = await kb_search(user_id, query, limit=KB_SEARCH_LIMIT, collection_ids=kb_cids)
                chunks = _filter_usable_chunks(chunks)
                if chunks and _kb_results_relevant(chunks, len(keywords)):
                    kb_chunks = [{"filename": c["filename"], "content": c["content"], "chunk_index": i}
                                 for i, c in enumerate(chunks)]
                    refs = "\n".join(f"【参考 {i+1} · 来源：{c['filename']}】\n{c['content'][:800]}"
                                    for i, c in enumerate(chunks))
                    system_prompt = system_prompt + "\n\n【知识库参考资料】\n" + refs + \
                        "\n\n以上资料根据关键词自动检索，仅供参考。请只关注用户本次发送的最新消息——只有当资料与本条消息直接相关时才引用。不要因为对话历史里聊过相关话题就引用资料，历史话题与当前问题可能已无关。引用格式：【参考 编号 · 来源：文件名】。如果不相关，直接忽略。"
        except Exception:
            logger.warning(f"KB search skipped in _stream")
    import json as _json
    full = ""
    try:
        _check_keywords(user_prompt, request)
        user_prompt = _mask_sensitive(user_prompt)
        async for chunk in call_llm_stream(system_prompt, user_prompt, api_key=api_key, model=model,
                                            temperature=temperature, max_tokens=max_tokens):
            if isinstance(chunk, dict):
                if "__done__" in chunk:
                    usage = chunk.get("usage", {})
                    tokens_val = usage.get("total_tokens", 0) if isinstance(usage, dict) else 0
                    full = _mask_sensitive(full)  # 输出脱敏
                    record_id = await add_history(user_id, type_, input_text, full, tokens=tokens_val)
                    _kb_chunks_out = _filter_referenced_chunks(kb_chunks, full)
                    yield f"data: {_json.dumps({'type':'done','usage':usage,'record_id':record_id,'kb_chunks':_kb_chunks_out}, ensure_ascii=False)}\n\n"
                elif "__error__" in chunk:
                    yield f"data: {_json.dumps({'type':'error','message':chunk['__error__']}, ensure_ascii=False)}\n\n"
            else:
                full += chunk
                yield f"data: {_json.dumps({'type':'token','content':chunk}, ensure_ascii=False)}\n\n"
    except Exception as e:
        logger.exception(f"流式生成异常 type={type_}")
        yield f"data: {_json.dumps({'type':'error','message':'服务内部错误，请稍后重试'}, ensure_ascii=False)}\n\n"


@app.post("/api/v1/summarize")
async def summarize(req: SummarizeRequest, request: Request, user: dict = Depends(get_current_user)):
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    system, user_p = build_summarize_prompt(req.text, req.length)

    if request.query_params.get("stream") == "true":
        return StreamingResponse(
            _stream("summarize", req.text[:2000], system, user_p, request, user["id"], api_key, model, custom_instruction=custom),
            media_type="text/event-stream",
            headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"}
        )

    data = await _call_and_save("summarize", req.text[:2000], system, user_p, request, user["id"], api_key, model, custom_instruction=custom)
    return _success(data)

@app.post("/api/v1/write-email")
async def write_email(req: EmailRequest, request: Request, user: dict = Depends(get_current_user)):
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    style = _get_style(request)
    temp = _style_temperature(style)
    extra = _style_instruction(style)
    system, user_p = build_email_prompt(req.recipient, req.subject_keywords, req.points, req.tone, req.original_email)
    system += extra
    if request.query_params.get("stream") == "true":
        return StreamingResponse(_stream("write-email", f"{req.recipient}:{req.subject_keywords}", system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})
    data = await _call_and_save("write-email", f"{req.recipient}:{req.subject_keywords}", system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom)
    return _success(data)

@app.post("/api/v1/meeting-minutes")
async def meeting_minutes(req: MinutesRequest, request: Request, user: dict = Depends(get_current_user)):
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    style = _get_style(request)
    temp = _style_temperature(style)
    extra = _style_instruction(style)
    system, user_p = build_minutes_prompt(req.transcript, req.speaker_tags)
    system += extra
    if request.query_params.get("stream") == "true":
        return StreamingResponse(_stream("meeting-minutes", req.transcript[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})
    data = await _call_and_save("meeting-minutes", req.transcript[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom)
    return _success(data)

@app.post("/api/v1/polish-report")
async def polish_report(req: PolishRequest, request: Request, user: dict = Depends(get_current_user)):
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    style = _get_style(request)
    temp = _style_temperature(style)
    extra = _style_instruction(style)
    system, user_p = build_polish_prompt(req.draft, req.style)
    system += extra
    if request.query_params.get("stream") == "true":
        return StreamingResponse(_stream("polish-report", req.draft[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})
    data = await _call_and_save("polish-report", req.draft[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom)
    return _success(data)

@app.post("/api/v1/report-ese")
async def report_ese(req: ReportEseRequest, request: Request, user: dict = Depends(get_current_user)):
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    style = _get_style(request)
    temp = _style_temperature(style)
    extra = _style_instruction(style)
    system, user_p = build_report_ese_prompt(req.rant, req.style)
    system += extra
    if request.query_params.get("stream") == "true":
        return StreamingResponse(_stream("report-ese", req.rant[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})
    data = await _call_and_save("report-ese", req.rant[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom)
    return _success(data)


@app.post("/api/v1/requirements")
async def requirements(req: RequirementsRequest, request: Request, user: dict = Depends(get_current_user)):
    """需求炼金 - 业务需求智能梳理"""
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    style = _get_style(request)
    temp = _style_temperature(style)
    extra = _style_instruction(style)
    system, user_p = build_requirements_prompt(req.text, req.style)
    system += extra
    if request.query_params.get("stream") == "true":
        return StreamingResponse(_stream("requirements", req.text[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})
    data = await _call_and_save("requirements", req.text[:2000], system, user_p, request, user["id"], api_key, model, temperature=temp, custom_instruction=custom)
    return _success(data)


@app.post("/api/v1/ppt-outline")
async def ppt_outline(req: PptOutlineRequest, request: Request, user: dict = Depends(get_current_user)):
    """PPT雕花 - PPT 大纲生成"""
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)
    system, user_p = build_ppt_outline_prompt(req.topic, req.points, req.style)
    if request.query_params.get("stream") == "true":
        return StreamingResponse(_stream("ppt-outline", f"{req.topic}:{req.points[:100]}", system, user_p, request, user["id"], api_key, model, custom_instruction=custom), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})
    data = await _call_and_save("ppt-outline", f"{req.topic}:{req.points[:100]}", system, user_p, request, user["id"], api_key, model, custom_instruction=custom)
    return _success(data)


@app.post("/api/v1/prd")
async def prd(req: PrdRequest, request: Request, user: dict = Depends(get_current_user)):
    """产品画饼 - 产品 PRD 生成（可选 Demo 由第二步单独生成）"""
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)

    # 第一步：生成 PRD（支持流式）
    system, user_p = build_prd_prompt(req.idea, req.style, False)

    # 流式模式
    if request.query_params.get("stream") == "true":
        async def prd_stream():
            nonlocal system
            import json as _json
            _check_keywords(user_p, request)
            user_p_masked = _mask_sensitive(user_p)
            # 知识库检索注入
            kb_chunks = []
            kb_col2 = request.headers.get("X-Kb-Collection", "")
            kb_cids = _parse_kb_collections(kb_col2) if kb_col2 else None
            if kb_cids and _need_kb_search(user_p):
                try:
                    keywords = await _extract_kb_keywords(user_p, api_key, model)
                    if keywords and _should_search_kb(keywords):
                        from database import kb_search
                        query = " OR ".join(keywords)
                        chunks = await kb_search(user["id"], query, limit=KB_SEARCH_LIMIT, collection_ids=kb_cids)
                        chunks = _filter_usable_chunks(chunks)
                        if chunks and _kb_results_relevant(chunks, len(keywords)):
                            kb_chunks = [{"filename": c["filename"], "content": c["content"], "chunk_index": i}
                                         for i, c in enumerate(chunks)]
                            refs = "\n".join(f"【参考 {i+1} · 来源：{c['filename']}】\n{c['content'][:800]}"
                                            for i, c in enumerate(chunks))
                            system = system + "\n\n【知识库参考资料】\n" + refs + \
                                "\n\n以上资料根据关键词自动检索，仅供参考。请只关注用户本次发送的最新消息——只有当资料与本条消息直接相关时才引用。不要因为对话历史里聊过相关话题就引用资料，历史话题与当前问题可能已无关。引用格式：【参考 编号 · 来源：文件名】。如果不相关，直接忽略。"
                except Exception:
                    logger.warning(f"KB search skipped in prd_stream")
            # 模板自定义指令注入
            if custom:
                system = system + "\n\n【用户自定义指令】\n" + custom
            full = ""
            try:
                async for chunk in call_llm_stream(system, user_p_masked, api_key=api_key, model=model, max_tokens=20000):
                    if isinstance(chunk, dict):
                        if "__done__" in chunk:
                            usage2 = chunk.get("usage", {})
                            tokens_val2 = usage2.get("total_tokens", 0) if isinstance(usage2, dict) else 0
                            full_masked = _mask_sensitive(full)
                            record_id = await add_history(user["id"], "prd", req.idea[:2000], full_masked, tokens=tokens_val2)
                            # PRD 完成，先通知前端
                            _kb_chunks_out = _filter_referenced_chunks(kb_chunks, full)
                            yield f"data: {_json.dumps({'type':'done','usage':usage2,'record_id':record_id,'kb_chunks':_kb_chunks_out}, ensure_ascii=False)}\n\n"
                            # 第二步：Demo（流式，完成后可编辑）
                            if req.with_demo:
                                try:
                                    demo_s, demo_u = build_prd_demo_prompt(full)
                                    demo_raw = ""
                                    demo_usage_val = {}
                                    async for dchunk in call_llm_stream(demo_s, demo_u, api_key=api_key, model=model, temperature=0.3, max_tokens=20000):
                                        if isinstance(dchunk, dict):
                                            if "__done__" in dchunk:
                                                demo_usage_val = dchunk.get("usage", {})
                                                demo_html = _strip_code_fence(demo_raw)
                                                if isinstance(usage2, dict) and isinstance(demo_usage_val, dict):
                                                    usage2["prompt_tokens"] = usage2.get("prompt_tokens", 0) + demo_usage_val.get("prompt_tokens", 0)
                                                    usage2["completion_tokens"] = usage2.get("completion_tokens", 0) + demo_usage_val.get("completion_tokens", 0)
                                                    usage2["total_tokens"] = usage2.get("total_tokens", 0) + demo_usage_val.get("total_tokens", 0)
                                                yield f"data: {_json.dumps({'type':'demo_done','demo_html':demo_html,'usage':usage2}, ensure_ascii=False)}\n\n"
                                            elif "__error__" in dchunk:
                                                yield f"data: {_json.dumps({'type':'demo_error','message':dchunk.get('__error__','Demo生成失败')}, ensure_ascii=False)}\n\n"
                                        else:
                                            demo_raw += dchunk
                                            yield f"data: {_json.dumps({'type':'demo_token','content':dchunk}, ensure_ascii=False)}\n\n"
                                except Exception as e:
                                    logger.warning(f"Demo生成失败: {e}", extra={"request_id": getattr(request.state, "request_id", "-")})
                                    yield f"data: {_json.dumps({'type':'demo_error','message':str(e)}, ensure_ascii=False)}\n\n"
                        elif "__error__" in chunk:
                            yield f"data: {_json.dumps({'type':'error','message':chunk['__error__']}, ensure_ascii=False)}\n\n"
                    else:
                        full += chunk
                        yield f"data: {_json.dumps({'type':'token','content':chunk}, ensure_ascii=False)}\n\n"
            except Exception as e:
                logger.exception(f"PRD流式生成异常")
                yield f"data: {_json.dumps({'type':'error','message':'服务内部错误，请稍后重试'}, ensure_ascii=False)}\n\n"
        return StreamingResponse(prd_stream(), media_type="text/event-stream", headers={"X-Accel-Buffering":"no","Cache-Control":"no-cache"})

    # 非流式模式
    data = await _call_and_save("prd", req.idea[:2000], system, user_p, request, user["id"], api_key, model, max_tokens=20000, custom_instruction=custom)

    demo_html = ""
    if req.with_demo:
        try:
            demo_system, demo_user = build_prd_demo_prompt(data["result"])
            demo_result, demo_usage = await call_llm(demo_system, demo_user, api_key=api_key, model=model, temperature=0.3, max_tokens=20000)
            demo_html = _strip_code_fence(demo_result)
            data["usage"]["prompt_tokens"] += demo_usage.get("prompt_tokens", 0)
            data["usage"]["completion_tokens"] += demo_usage.get("completion_tokens", 0)
            data["usage"]["total_tokens"] += demo_usage.get("total_tokens", 0)
        except Exception as e:
            logger.warning(f"Demo生成失败: {str(e)}", extra={"request_id": getattr(request.state, "request_id", "-")})

    data["demo_html"] = demo_html
    return _success(data)


@app.post("/api/v1/weekly-report")
async def weekly_report(req: WeeklyReportRequest, request: Request, user: dict = Depends(get_current_user)):
    """周报生成 — 四种模板 × 两种风格 × 中英双语"""
    await _check_rate_limit(user["id"], request.client.host, request)
    api_key, model, custom, _ = _get_headers(request)

    # 根据 report_type 选择对应的 prompt builder
    builders = {
        "research": build_weekly_research_prompt,
        "project": build_weekly_project_prompt,
        "techsurvey": build_weekly_techsurvey_prompt,
        "ops": build_weekly_ops_prompt,
    }

    if req.report_type not in builders:
        return _error(400, _t("invalid_report_type", request))

    builder = builders[req.report_type]
    system, user_p = builder(req.raw_notes, style=req.style, lang=req.lang)

    # 模板自定义指令注入
    if custom:
        system = system + "\n\n【用户自定义指令】\n" + custom

    input_preview = f"type={req.report_type} style={req.style} lang={req.lang}"
    if request.query_params.get("stream") == "true":
        return StreamingResponse(
            _stream("weekly-report", input_preview[:200], system, user_p, request, user["id"], api_key, model, custom_instruction=custom),
            media_type="text/event-stream",
            headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"}
        )
    data = await _call_and_save("weekly-report", input_preview[:200], system, user_p, request, user["id"], api_key, model, custom_instruction=custom)
    return _success(data)


# ============================================================
#  余额查询（需登录）
# ============================================================

@app.get("/api/v1/balance")
async def get_balance(request: Request, user: dict = Depends(get_current_user)):
    import httpx
    api_key, model, custom, _ = _get_headers(request)
    if not api_key:
        return _error(400, _t("api_key_required", request))
    try:
        async with httpx.AsyncClient(timeout=8.0) as client:
            resp = await client.get(
                "https://api.deepseek.com/user/balance",
                headers={"Authorization": f"Bearer {api_key}"},
            )
            if resp.status_code == 200:
                return _success(resp.json())
            else:
                return _error(resp.status_code, _t("balance_query_failed", request, detail=resp.text[:200]))
    except Exception as e:
        return _safe_error(502, _t("balance_query_failed", request), e, request)


# ============================================================
#  历史记录（需登录，只能操作自己的）
# ============================================================

@app.get("/api/v1/stats")
async def get_stats(user: dict = Depends(get_current_user)):
    """用户使用统计"""
    records = await get_history(user["id"], limit=1000)
    total = len(records)
    types = {}
    total_tokens = 0
    for r in records:
        t = r.get("type", "other")
        types[t] = types.get(t, 0) + 1
        total_tokens += r.get("tokens", 0) if isinstance(r.get("tokens"), int) else 0
    top_type = max(types, key=types.get) if types else "none"
    return _success({
        "total_generations": total,
        "top_feature": top_type,
        "total_tokens": total_tokens,
        "by_type": types,
        "estimated_chars": total_tokens * 2,  # 估算字数
    })

@app.get("/api/v1/history")
async def list_history(user: dict = Depends(get_current_user), limit: int = 50):
    """获取用户历史记录列表"""
    records = await get_history(user["id"], limit=min(limit, 1000))
    return _success({"records": records})

@app.delete("/api/v1/history/{record_id}")
async def remove_history(record_id: int, request: Request, user: dict = Depends(get_current_user)):
    ok = await delete_history(record_id, user["id"])
    if not ok:
        return _error(404, _t("record_not_found_or_denied", request))
    return _success(None, _t("record_deleted", request))


@app.get("/api/v1/history/{record_id}")
async def fetch_history(record_id: int, request: Request, user: dict = Depends(get_current_user)):
    """获取单条历史记录完整内容"""
    record = await get_history_by_id(record_id, user["id"])
    if not record:
        return _error(404, _t("record_not_found_or_denied", request))
    return _success({"id": record["id"], "type": record["type"], "input_text": record["input_text"],
                     "result_text": record["result_text"], "tokens": record.get("tokens", 0),
                     "created_at": record["created_at"], "kb_chunks": []})


# ============================================================
#  健康检查（公开）
# ============================================================

@app.post("/api/v1/export-docx")
async def export_docx(request: Request, user: dict = Depends(get_current_user)):
    """Markdown 内容导出为 .docx 文件"""
    import io
    try:
        import docx
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return _error(501, "需要 pip install python-docx")

    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("correct_json", request))

    content = str(body.get("content", ""))
    title = str(body.get("title", "export"))[:50]
    lang = str(body.get("lang", "zh"))

    try:
        doc = docx.Document()
        from docx.oxml.ns import qn as _qn

        def _set_cjk_font(run, font_name):
            """设置东亚字体（python-docx 的 font.name 只设西文，CJK 需额外操作 XML）"""
            run.font.name = font_name
            rPr = run._element.get_or_add_rPr()
            rFonts = rPr.find(_qn('w:rFonts'))
            if rFonts is None:
                from lxml import etree
                rFonts = etree.SubElement(rPr, _qn('w:rFonts'))
            rFonts.set(_qn('w:eastAsia'), font_name)
            rFonts.set(_qn('w:ascii'), font_name)
            rFonts.set(_qn('w:hAnsi'), font_name)

        header = doc.sections[0].header
        h = header.add_paragraph("GrindPal · AI生成内容")
        h.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        for run in h.runs:
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(150, 150, 150)

        doc.add_heading(title, level=1)

        # Markdown 逐行解析
        import re as _md_re
        in_code_block = False
        code_buffer = []
        ordered_list_counter = 0
        in_table = False
        table_rows = []  # 暂存表格行数据

        def _add_run(para, text, bold=False, italic=False, code=False, font_size=None, zh_font=False):
            """向段落添加一个 run，支持富文本"""
            run = para.add_run(text)
            if bold:
                run.bold = True
            if italic:
                run.italic = True
            if code:
                run.font.name = 'Consolas'
                run.font.size = font_size or Pt(9.5)
                run.font.color.rgb = RGBColor(0xE8, 0x3E, 0x8C)
            elif zh_font:
                _set_cjk_font(run, '仿宋')
            return run

        def _parse_inline(text):
            """解析行内格式：**bold**, *italic*, `code`, [link](url)"""
            parts = []
            pos = 0
            for match in _md_re.finditer(r'(\*\*(.+?)\*\*)|(\*(.+?)\*)|(`[^`]+`)|(\[([^\]]+)\]\(([^)]+)\))', text):
                if match.start() > pos:
                    parts.append(('text', text[pos:match.start()]))
                if match.group(2):  # **bold**
                    parts.append(('bold', match.group(2)))
                elif match.group(4):  # *italic*
                    parts.append(('italic', match.group(4)))
                elif match.group(0).startswith('`'):  # `code`
                    parts.append(('code', match.group(0)[1:-1]))
                elif match.group(6):  # [link](url)
                    parts.append(('text', f'{match.group(6)} ({match.group(7)})'))
                pos = match.end()
            if pos < len(text):
                parts.append(('text', text[pos:]))
            return parts

        def _parse_table_row(line):
            """解析表格行: | col1 | col2 | col3 | → [col1, col2, col3]"""
            cells = line.strip().strip('|').split('|')
            return [c.strip() for c in cells]

        def _is_table_sep(line):
            """判断是否为表格分隔行: |---|---|"""
            return bool(_md_re.match(r'^\|?[\s]*:?---+:?[\s]*(\|[\s]*:?---+:?[\s]*)+\|?$', line))

        def _flush_table():
            """将暂存的表格行写入 doc"""
            nonlocal table_rows
            if not table_rows:
                return
            rows = len(table_rows)
            cols = max(len(r) for r in table_rows) if table_rows else 1
            table = doc.add_table(rows=rows, cols=cols, style='Table Grid')
            for ri, row_data in enumerate(table_rows):
                for ci, cell_text in enumerate(row_data):
                    if ci < cols:
                        cell = table.cell(ri, ci)
                        cell.text = ''
                        p = cell.paragraphs[0]
                        for kind, val in _parse_inline(cell_text):
                            _add_run(p, val, bold=(kind == 'bold'), italic=(kind == 'italic'),
                                     code=(kind == 'code'), zh_font=(lang == 'zh'))
            doc.add_paragraph()  # 表格后空一行
            table_rows = []

        for line in content.split('\n'):
            stripped = line.strip()

            # 表格处理
            if _is_table_sep(stripped):
                in_table = True
                continue

            if stripped.startswith('|') and (stripped.endswith('|') or '|' in stripped[1:]):
                if in_table:
                    table_rows.append(_parse_table_row(stripped))
                    continue
                elif table_rows:
                    # 追加行（多行表格）
                    table_rows.append(_parse_table_row(stripped))
                    continue
                else:
                    # 新表格开始：第一行是表头，需要等分隔行确认
                    table_rows = [_parse_table_row(stripped)]
                    continue
            else:
                if in_table and table_rows:
                    _flush_table()
                in_table = False

            # 代码块切换
            if stripped.startswith('```') or stripped.startswith('~~~'):
                if in_code_block:
                    for cb_line in code_buffer:
                        p = doc.add_paragraph()
                        p.paragraph_format.keep_with_next = True
                        run = p.add_run(cb_line)
                        run.font.name = 'Consolas'
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
                        p.paragraph_format.space_before = Pt(0)
                        p.paragraph_format.space_after = Pt(0)
                        p.paragraph_format.left_indent = Cm(0.5)
                    code_buffer = []
                    in_code_block = False
                else:
                    in_code_block = True
                    code_buffer = []
                continue

            if in_code_block:
                code_buffer.append(stripped)
                continue

            if not stripped:
                doc.add_paragraph()
                ordered_list_counter = 0
                continue

            # 标题
            if stripped.startswith('#### '):
                p = doc.add_heading(stripped[5:], level=4)
                if lang == 'zh':
                    for run in p.runs:
                        _set_cjk_font(run, '仿宋')
                continue
            elif stripped.startswith('### '):
                p = doc.add_heading(stripped[4:], level=3)
                if lang == 'zh':
                    for run in p.runs:
                        _set_cjk_font(run, '仿宋')
                continue
            elif stripped.startswith('## '):
                p = doc.add_heading(stripped[3:], level=2)
                if lang == 'zh':
                    for run in p.runs:
                        _set_cjk_font(run, '仿宋')
                continue
            elif stripped.startswith('# '):
                p = doc.add_heading(stripped[2:], level=2)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if lang == 'zh':
                    for run in p.runs:
                        _set_cjk_font(run, '仿宋')
                continue

            # 有序列表
            ol_match = _md_re.match(r'^(\d+)\.\s+(.*)', stripped)
            if ol_match:
                ordered_list_counter = int(ol_match.group(1))
                p = doc.add_paragraph(style='List Number')
                p.clear()
                for kind, val in _parse_inline(ol_match.group(2)):
                    _add_run(p, val, bold=(kind=='bold'), italic=(kind=='italic'), code=(kind=='code'), zh_font=(lang=='zh'))
                if lang == 'zh':
                    for run in p.runs:
                        if run.font.name != 'Consolas':
                            _set_cjk_font(run, '仿宋')
                continue

            # 无序列表
            if stripped.startswith('- ') or stripped.startswith('* '):
                item_text = stripped[2:]
                p = doc.add_paragraph(style='List Bullet')
                p.clear()
                for kind, val in _parse_inline(item_text):
                    _add_run(p, val, bold=(kind=='bold'), italic=(kind=='italic'), code=(kind=='code'), zh_font=(lang=='zh'))
                if lang == 'zh':
                    for run in p.runs:
                        if run.font.name != 'Consolas':
                            _set_cjk_font(run, '仿宋')
                ordered_list_counter = 0
                continue

            # 引用块
            if stripped.startswith('> '):
                quote_text = stripped[2:]
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = Cm(1)
                p.paragraph_format.space_before = Pt(2)
                p.paragraph_format.space_after = Pt(2)
                run = p.add_run(quote_text)
                run.font.italic = True
                run.font.color.rgb = RGBColor(100, 100, 100)
                if lang == 'zh':
                    _set_cjk_font(run, '仿宋')
                ordered_list_counter = 0
                continue

            # 水平分割线
            if stripped in ('---', '***', '___'):
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(6)
                run = p.add_run('─' * 60)
                run.font.color.rgb = RGBColor(200, 200, 200)
                run.font.size = Pt(8)
                ordered_list_counter = 0
                continue

            # 普通段落（含行内格式）
            ordered_list_counter = 0
            p = doc.add_paragraph()
            for kind, val in _parse_inline(stripped):
                _add_run(p, val, bold=(kind=='bold'), italic=(kind=='italic'), code=(kind=='code'), zh_font=(lang=='zh'))
            if lang == 'zh':
                for run in p.runs:
                    if run.font.name != 'Consolas':
                        _set_cjk_font(run, '仿宋')
                p.paragraph_format.first_line_indent = Cm(0.74)

        # 处理末尾未刷新的表格
        if in_table and table_rows:
            _flush_table()

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
    except Exception as e:
        return _safe_error(500, _t("generate_failed", request), e, request)

    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f"attachment; filename=export.docx"})


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """#1E3A5F → (30, 58, 95)"""
    h = hex_color.lstrip('#')
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


@app.post("/api/v1/export-pptx")
async def export_pptx(request: Request, user: dict = Depends(get_current_user)):
    """PPT 大纲导出为 .pptx 文件"""
    import io
    try:
        import pptx
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return _error(501, _t("pptx_need_install", request))

    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("correct_json", request))

    outline_id = body.get("outline_id")
    content = body.get("content", "")
    theme = body.get("theme", "blue")

    if outline_id:
        record = await get_history_by_id(int(outline_id), user["id"])
        if not record:
            return _error(404, _t("outline_not_found", request))
        if record.get("type") != "ppt-outline":
            return _error(400, _t("outline_not_ppt", request))
        content = str(record.get("result_text", ""))
    if not content:
        return _error(400, _t("outline_need_params", request))

    themes = {
        "blue":  ("#1E3A5F", "#3B82F6", "#D0D0D0"),
        "gray":  ("#2D3436", "#636E72", "#D0D0D0"),
        "warm":  ("#C2410C", "#F59E0B", "#D0D0D0"),
    }
    title_color, bullet_color, footer_color = themes.get(theme, themes["blue"])

    try:
        prs = pptx.Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        # 查找空白布局（按名称而非索引，兼容不同版本）
        blank_layout = None
        for layout in prs.slide_layouts:
            if layout.name == "Blank":
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = prs.slide_layouts[len(prs.slide_layouts) - 1]  # 兜底取最后一个

        lines = content.strip().split('\n')
        slides_data = []
        current = {"title": "", "bullets": [], "notes": ""}
        in_notes = False
        for line in lines:
            s = line.strip()
            if not s:
                if current["title"]:
                    slides_data.append(current)
                    current = {"title": "", "bullets": [], "notes": ""}
                in_notes = False
                continue
            if s.startswith('## ') or s.startswith('# '):
                if current["title"]:
                    slides_data.append(current)
                    current = {"title": "", "bullets": [], "notes": ""}
                current["title"] = s.lstrip('#').strip()
                in_notes = False
            elif s.startswith('- ') or s.startswith('* '):
                current["bullets"].append(s.lstrip('-* ').strip())
                in_notes = False
            elif s.startswith('备注') or s.startswith('>'):
                in_notes = True
                current["notes"] = s.lstrip('备注>：: ').strip()
            elif in_notes:
                current["notes"] += '\n' + s
        if current["title"]:
            slides_data.append(current)
        if not slides_data:
            return _error(400, _t("outline_parse_failed", request))

        total = len(slides_data)
        for i, sd in enumerate(slides_data, 1):
            if not sd["title"]:
                continue
            slide = prs.slides.add_slide(blank_layout)
            # 标题
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = sd["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*_hex_to_rgb(title_color))
            # 分隔线
            line_shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(11.7), Pt(2))
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(bullet_color))
            line_shape.line.fill.background()
            # 要点
            tb2 = slide.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.5), Inches(4.5))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            for j, b in enumerate(sd["bullets"]):
                if j > 0:
                    tf2.add_paragraph()
                tf2.paragraphs[j].text = b
                tf2.paragraphs[j].font.size = Pt(20)
                tf2.paragraphs[j].font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                tf2.paragraphs[j].space_after = Pt(8)
            # 备注
            if sd["notes"].strip():
                ns = slide.notes_slide
                ns.notes_text_frame.text = sd["notes"]
            # 页码 + 水印
            fb = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
            fb.text_frame.paragraphs[0].text = f"{i}/{total}"
            fb.text_frame.paragraphs[0].font.size = Pt(10)
            fb.text_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(footer_color))
            fb.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            wm = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(3), Inches(0.4))
            wm.text_frame.paragraphs[0].text = "GrindPal · AI 生成"
            wm.text_frame.paragraphs[0].font.size = Pt(9)
            wm.text_frame.paragraphs[0].font.color.rgb = RGBColor(*_hex_to_rgb(footer_color))

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
    except Exception as e:
        return _safe_error(500, _t("generate_failed", request), e, request)

    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=outline.pptx"})


@app.post("/api/v1/transcribe")
async def transcribe_audio(request: Request, user: dict = Depends(get_current_user)):
    """上传音频文件转文字（openai-whisper base，稳定批量处理）"""
    global _whisper_model, _whisper_lock
    if _whisper_lock is None:
        _whisper_lock = asyncio.Lock()
    async with _whisper_lock:
        if _whisper_model is None:
            try:
                import whisper, asyncio as _asyncio
                _whisper_model = await _asyncio.to_thread(whisper.load_model, "base")
            except ImportError:
                return _error(501, _t("whisper_need_install", request))
            except Exception as e:
                return _safe_error(501, _t("model_load_failed", request), e, request)

    if size_err := _check_upload_size(request): return size_err
    form = await request.form()
    file = form.get("file")
    if not file:
        return _error(400, _t("audio_no_file", request))

    import tempfile, os
    suffix = os.path.splitext(file.filename or "audio.wav")[1] or ".wav"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        result = _whisper_model.transcribe(tmp_path, language="zh")
        text = result["text"].strip()
        return _success({"text": text, "filename": file.filename})
    except Exception as e:
        return _safe_error(500, "transcribe_failed", e, request)
    finally:
        os.unlink(tmp_path)


def _text_similar(a: str, b: str) -> bool:
    """判断两段文本是否高度重复 — 多维度综合判断"""
    if not a or not b:
        return False
    # 某一条太短（<3字）直接判相似
    if len(a) < 3 or len(b) < 3:
        return True
    # 短文本 Jaccard 阈值更严格
    threshold = 0.55 if min(len(a), len(b)) < 10 else 0.65
    set_a, set_b = set(a), set(b)
    if not set_a or not set_b:
        return False
    intersection = len(set_a & set_b)
    union = len(set_a | set_b)
    jaccard = intersection / union if union > 0 else 0
    # b 基本包含 a 的全部字符
    contained = len(set_a - set_b) <= max(2, len(set_a) * 0.15)
    return jaccard > threshold or contained

def _has_speech(wav_path: str) -> bool:
    """用 webrtcvad + RMS 能量双重检测 wav 中是否有人声"""
    import wave, struct, webrtcvad
    vad = webrtcvad.Vad(3)
    try:
        with wave.open(wav_path, 'rb') as wf:
            if wf.getnchannels() != 1 or wf.getsampwidth() != 2 or wf.getframerate() != 16000:
                return True
            raw = wf.readframes(wf.getnframes())
        samples = struct.unpack(f"<{len(raw)//2}h", raw)
        max_abs = max(abs(s) for s in samples) if samples else 0
        if max_abs < 200:
            return False
        frame_len = 480 * 2
        speech_frames = sum(
            1 for i in range(0, len(raw) - frame_len + 1, frame_len)
            if vad.is_speech(raw[i:i+frame_len], 16000)
        )
        total_frames = max(len(raw) // frame_len, 1)
        return (speech_frames / total_frames) > 0.2
    except Exception:
        logger.warning("VAD 语音检测失败，默认有语音", exc_info=True)
        return True

def _load_whisper_model(model_name: str):
    """加载或切换 whisper.cpp 模型。返回 Model 实例。
    若 pywhispercpp 未安装，抛出 ImportError。"""
    global _wcp_model, _wcp_current_model_name, _n_gpu_layers
    import os as _os
    try:
        import pywhispercpp.constants as _pw_constants
    except ImportError:
        raise ImportError("pywhispercpp 未安装，语音转写功能不可用。请执行: pip install pywhispercpp webrtcvad")
    # 国内网络环境：HF 不可用时走清华镜像
    # 国内网络环境：HF 不可用时走清华镜像
    if not hasattr(_pw_constants, '_MIRROR_PATCHED'):
        _pw_constants.MODELS_BASE_URL = "https://hf-mirror.com"
        _pw_constants._MIRROR_PATCHED = True
    from pywhispercpp.model import Model

    models_dir = _os.path.expanduser("~/.local/share/pywhispercpp/models")
    ggml_path = _os.path.join(models_dir, f"ggml-{model_name}.bin")
    # 如果本地文件不存在，传模型名让 pywhispercpp 自动下载
    resolved = ggml_path if _os.path.exists(ggml_path) else model_name

    if _wcp_current_model_name != model_name or _wcp_model is None:
        kwargs = {}
        # pywhispercpp 1.5.0 不支持 n_gpu_layers 参数，GPU 由 whisper.cpp 自动管理
        # 如需指定 GPU 设备，通过 context_params 传递
        # if _n_gpu_layers > 0:
        #     kwargs["context_params"] = {"use_gpu": True, "gpu_device": 0}
        try:
            _wcp_model = Model(resolved, **kwargs)
            _wcp_current_model_name = model_name
        except Exception as e:
            # 下载失败（如网络不通）→ 回退到已有模型，优先大模型
            import glob as _glob, re as _re
            existing = sorted(_glob.glob(_os.path.join(models_dir, "ggml-*.bin")))
            if existing:
                # 按模型大小优先级排序：medium > small > base > tiny
                _order = {"medium": 0, "small": 1, "base": 2, "tiny": 3}
                existing.sort(key=lambda p: _order.get(
                    _re.sub(r'^ggml-|\.bin$', '', _os.path.basename(p)), 99
                ))
                fallback = _os.path.basename(existing[0])
                fallback_name = _re.sub(r'^ggml-|\.bin$', '', fallback)
                logger.warning(
                    f"模型 {model_name} 加载失败（{e}），回退到 {fallback_name}"
                )
                _wcp_model = Model(existing[0], **kwargs)
                _wcp_current_model_name = fallback_name
            else:
                raise  # 没有任何模型，抛出原始错误
    return _wcp_model


@app.post("/api/v1/transcribe-vad")
async def transcribe_vad(request: Request, user: dict = Depends(get_current_user)):
    """whisper.cpp 快速转写（支持 tiny/base/small/medium 模型选择）"""
    global _wcp_lock, _webm_header
    if _wcp_lock is None:
        _wcp_lock = asyncio.Lock()

    if size_err := _check_upload_size(request): return size_err
    form = await request.form()
    file = form.get("file")
    language = (form.get("language") or "zh")
    model_name = (form.get("model") or "base").strip()
    if model_name not in ("base", "small", "medium"):
        model_name = "base"
    gpu_layers = 0  # GPU加速已禁用

    try:
        import pywhispercpp
    except ImportError:
        return _error(501, "pywhispercpp 未安装，语音转写功能不可用。请执行: pip install pywhispercpp webrtcvad")

    import tempfile, asyncio as _asyncio, subprocess

    suffix = os.path.splitext(file.filename or "audio.wav")[1] or ".webm"
    raw_data = await file.read()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(raw_data)
        tmp_path = tmp.name
    wav_path = tmp_path + ".wav"

    try:
        # 模型加载（加锁保护，支持热切换；GPU 层数变化时重新加载）
        async with _wcp_lock:
            global _n_gpu_layers
            try:
                _wcp_model = await _asyncio.to_thread(lambda: _load_whisper_model(model_name))
            except ImportError as e:
                return _error(501, str(e))
            # GPU加速已禁用
            # if gpu_layers != _n_gpu_layers:
            #     _n_gpu_layers = gpu_layers
            #     _wcp_current_model_name = None
            _wcp_model = await _asyncio.to_thread(lambda: _load_whisper_model(model_name))
        # 手动 ffmpeg 转换
        def _run_ffmpeg(src_path):
            subprocess.run(
                ["ffmpeg", "-err_detect", "ignore_err", "-i", src_path,
                 "-af", "apad=whole_dur=0.5",
                 "-ar", "16000", "-ac", "1", "-c:a", "pcm_s16le", wav_path, "-y"],
                capture_output=True, check=True,
            )

        try:
            await _asyncio.to_thread(lambda: _run_ffmpeg(tmp_path))
        except subprocess.CalledProcessError as e:
            stderr_text = e.stderr.decode(errors='replace')
            # Chrome MediaRecorder 后续 chunk 缺少 EBML 头 → 补头重试
            if "EBML header" in stderr_text and _webm_header is not None:
                with open(tmp_path, "wb") as f:
                    f.write(_webm_header + raw_data)
                await _asyncio.to_thread(lambda: _run_ffmpeg(tmp_path))
            else:
                raise

        # 首次成功 → 保存 EBML 头供后续 chunk 使用
        # 取第一个 Cluster 之前的所有字节作为容器头，避免混入音频数据
        if _webm_header is None:
            cluster_marker = b'\x1F\x43\xB6\x75'
            idx = raw_data.find(cluster_marker)
            _webm_header = raw_data[:idx] if idx > 0 else raw_data[:800]

        # VAD 静音检测：没人说话直接返回空，不浪费 whisper 算力
        if not _has_speech(wav_path):
            return _success({"text": "", "filename": file.filename, "silent": True})

        # 用上一次转写文本做 initial_prompt，提供上下文
        async with _last_text_lock:
            prev = _last_text.get(user["id"], "")
            # 超过 60 秒无更新 → 清上下文，可能是新一段录音
            import time as _time
            _last_ts = _last_text.get(f"_{user['id']}_ts", 0)
            if _time.time() - _last_ts > 60:
                prev = ""
        prompt = prev if prev else ("以下是以普通话为主的对话，里面可能夹杂有英文或空白，请如实输出你识别的内容。" if language == "zh" else "")
        async with _wcp_lock:
            _wcp_model._params.language = language
            segments = await _asyncio.to_thread(
                lambda: _wcp_model.transcribe(wav_path, initial_prompt=prompt, temperature=0.0)
            )
        text = " ".join(s.text.strip() for s in segments)
        # 结果太短（<3字符）→ 大概率是幻觉，丢弃
        if len(text) < 3:
            return _success({"text": "", "filename": file.filename, "short": True})
        # 与上次结果相似度过高 → 重复，跳过
        async with _last_text_lock:
            if prev and _text_similar(text, prev):
                return _success({"text": "", "filename": file.filename, "dup": True})
            _last_text[user["id"]] = text
            _last_text[f"_{user['id']}_ts"] = _time.time()
            # 清理超过 5 分钟的旧条目
            cutoff = _time.time() - 300
            stale = [k for k in list(_last_text.keys()) if not str(k).startswith('_') and _last_text.get(f'_{k}_ts', 0) < cutoff]
            for k in stale:
                del _last_text[k]
                _last_text.pop(f"_{k}_ts", None)
        return _success({"text": text, "filename": file.filename})
    except subprocess.CalledProcessError as e:
        stderr_text = e.stderr.decode(errors='replace')
        # ffmpeg 真正报错在末尾，取最后 500 字符
        return _error(500, _t("audio_convert_failed", request, detail=stderr_text[-500:]))
    except Exception as e:
        return _safe_error(500, "transcribe_failed", e, request)
    finally:
        os.unlink(tmp_path)
        if os.path.exists(wav_path):
            os.unlink(wav_path)


@app.post("/api/v1/transcribe-batch")
async def transcribe_batch(request: Request, user: dict = Depends(get_current_user)):
    """录制完成后的全文转写：接收完整 webm，一次性转写全文"""
    global _wcp_lock
    if _wcp_lock is None:
        _wcp_lock = asyncio.Lock()

    if size_err := _check_upload_size(request): return size_err
    import tempfile, asyncio as _asyncio, subprocess
    form = await request.form()
    file = form.get("file")
    language = (form.get("language") or "zh")
    model_name = (form.get("model") or "medium").strip()
    if model_name not in ("base", "small", "medium"):
        model_name = "medium"
    gpu_layers = 0  # GPU加速已禁用

    try:
        import pywhispercpp
    except ImportError:
        return _error(501, "pywhispercpp 未安装，语音转写功能不可用。请执行: pip install pywhispercpp webrtcvad")

    suffix = os.path.splitext(file.filename or "audio.wav")[1] or ".webm"
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name
    wav_path = tmp_path + ".wav"
    try:
        # 模型加载（加锁保护，支持热切换；GPU 层数变化时重新加载）
        async with _wcp_lock:
            global _n_gpu_layers
            if gpu_layers != _n_gpu_layers:
                _n_gpu_layers = gpu_layers
                _wcp_current_model_name = None
            _wcp_model = await _asyncio.to_thread(lambda: _load_whisper_model(model_name))
        await _asyncio.to_thread(
            lambda: subprocess.run(
                ["ffmpeg", "-err_detect", "ignore_err", "-i", tmp_path,
                 "-af", "apad=whole_dur=0.5",
                 "-ar", "16000", "-ac", "1", "-c:a", "pcm_s16le", wav_path, "-y"],
                capture_output=True, check=True,
            )
        )
        prompt = "以下是以普通话为主的对话，里面可能夹杂有英文或空白，请如实输出你识别的内容。" if language == "zh" else ""
        async with _wcp_lock:
            _wcp_model._params.language = language
            segments = await _asyncio.to_thread(
                lambda: _wcp_model.transcribe(wav_path, initial_prompt=prompt, temperature=0.0)
            )
        text = " ".join(s.text.strip() for s in segments)
        return _success({"text": text, "filename": file.filename})
    except subprocess.CalledProcessError as e:
        stderr_text = e.stderr.decode(errors='replace')
        return _error(500, _t("audio_convert_failed", request, detail=stderr_text[-500:]))
    except Exception as e:
        return _safe_error(500, "transcribe_failed", e, request)
    finally:
        os.unlink(tmp_path)
        if os.path.exists(wav_path):
            os.unlink(wav_path)


@app.post("/api/v1/correct-text")
async def correct_text(request: Request, user: dict = Depends(get_current_user)):
    """LLM 纠错顺滑：修正语音识别错误，删除重复，理顺语序"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("correct_json", request))
    raw_text = str(body.get("text", "")).strip()
    if not raw_text:
        return _error(400, _t("correct_no_text", request))
    api_key, model, custom, _ = _get_headers(request)
    system_prompt = (
        "# Role: 语音识别校对员\n"
        "## Profile\n"
        "你修正语音转文字中的三类错误：同音字（「工能」→「功能」）、断句错误、重复内容。你只做修正，不改写。\n"
        "## Rules\n"
        "1. 只修正明显的识别错误，不确定的地方保留原文\n"
        "2. 同一句话多次出现只保留一次\n"
        "3. 补充必要的标点（逗号、句号、问号）\n"
        "4. 不改变原意、不添加新信息、不美化表达\n"
        "5. 只输出校对后文本"
    )
    if custom:
        system_prompt = system_prompt + "\n\n【用户自定义指令】\n" + custom
    try:
        content, _ = await call_llm(
            system_prompt, f"请校对以下语音识别结果：\n\n{raw_text}",
            api_key=api_key, model=model,
            temperature=0.3, max_tokens=2048,
        )
        corrected = (content or "").strip()
        if not corrected:
            # LLM 返回空文本，回退原文
            return _success({"text": raw_text, "fallback": True})
        return _success({"text": corrected})
    except Exception as e:
        return _safe_error(502, "correct_failed", e, request)


@app.get("/api/v1/health")
async def health():
    return _success({"status": "ok", "version": VERSION, "ocr_available": False})


# ---- 文件上传提取文本 ----
@app.post("/api/v1/extract-text")
async def extract_text(request: Request, user: dict = Depends(get_current_user)):
    """上传文件提取文本（支持 txt/md/docx/pdf/xlsx/xls/pptx/csv/json 等常见格式）"""
    content_type = request.headers.get("content-type", "")
    if "multipart/form-data" not in content_type:
        return _error(400, _t("file_upload_hint", request))

    form = await request.form()
    file = form.get("file")
    if not file:
        return _error(400, _t("file_missing", request))

    filename = file.filename or ""
    raw = await file.read()
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    try:
        text = _extract_doc_text(raw, ext, filename)
    except Exception as e:
        return _safe_error(500, _t("file_parse_failed", request), e)

    if len(text) > 50000:
        text = text[:50000] + "\n…(内容过长已截断)"
    return _success({"text": text, "filename": filename, "length": len(text)})


# ============================================================
#  [自由对话] 端点
# ============================================================

@app.post("/api/v1/chat/conversations")
async def chat_create_conversation(request: Request, user: dict = Depends(get_current_user)):
    """创建新对话"""
    try:
        body = await request.json()
    except Exception:
        body = {}
    title = str(body.get("title", "新对话")).strip() or "新对话"
    conv_id = await create_conversation(user["id"], title)
    if not conv_id:
        return _error(500, _t("conversation_create_failed", request))
    return _success({"id": conv_id, "title": title}, "conversation_created", request)


@app.get("/api/v1/chat/conversations")
async def chat_list_conversations(
    limit: int = 100,
    offset: int = 0,
    user: dict = Depends(get_current_user)
):
    """列出用户的所有对话，支持分页"""
    convs = await list_conversations(user["id"], limit=limit, offset=offset)
    return _success({"conversations": convs})


@app.get("/api/v1/chat/conversations/{conv_id}")
async def chat_get_conversation(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """获取单个对话"""
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    return _success({"conversation": conv})


@app.get("/api/v1/chat/search")
async def chat_search(
    q: str = "",
    limit: int = 20,
    request: Request = None,
    user: dict = Depends(get_current_user)
):
    """全文搜索用户的所有对话消息"""
    await _check_rate_limit(user["id"], request.client.host, request)
    q = (q or "").strip()
    if not q or len(q) < 2:
        return _success({"results": []})
    # 清洗 FTS5 查询语法
    safe_q = q.replace('"', '').replace("'", "").replace("*", "")
    results = await search_messages(user["id"], safe_q, limit=limit)
    return _success({"results": results})


@app.get("/api/v1/chat/conversations/{conv_id}/context-usage")
async def chat_context_usage(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """估算当前对话的 token 使用量"""
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    msgs = await get_messages(conv_id)
    # 简单估算：中文 ~1.5 char/token, 英文 ~4 char/token，取中间 ~2.5
    total_chars = sum(len(m["content"] or "") for m in msgs)
    # 加上 system prompt (~600 chars) 和 SAFETY_GUARD (~400 chars)
    total_chars += 1000
    estimated_tokens = max(1, int(total_chars / 2.5))
    # 默认模型上下文 128K
    model_name = request.headers.get("X-Model", "") or os.getenv("DEEPSEEK_MODEL", "deepseek-chat")
    model_max = 128000
    if "reasoner" in model_name:
        model_max = 128000
    elif "flash" in model_name:
        model_max = 128000
    usage_percent = round(estimated_tokens / model_max * 100, 1) if model_max else 0
    return _success({
        "estimated_tokens": estimated_tokens,
        "model_max": model_max,
        "usage_percent": min(usage_percent, 100),
        "message_count": len(msgs),
    })


@app.get("/api/v1/chat/conversations/{conv_id}/export")
async def chat_export_conversation(conv_id: int, request: Request, format: str = "md", user: dict = Depends(get_current_user)):
    """导出对话为 Markdown 或纯文本"""
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    msgs = await get_messages(conv_id, limit=5000)
    title = conv.get("title", "新对话")
    created = conv.get("created_at", "")
    lines = [f"# {title}", f"", f"> 导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}", f"> 创建时间：{created}", f"> 消息数：{len(msgs)}", f""]
    for m in msgs:
        role_label = "👤 用户" if m["role"] == "user" else "🤖 AI"
        time_str = (m.get("created_at") or "")[:19]
        lines.append(f"### {role_label} ({time_str})")
        lines.append("")
        lines.append(m["content"])
        lines.append("")
        lines.append("---")
        lines.append("")
    md_content = "\n".join(lines)
    from urllib.parse import quote as _url_quote
    safe_title = title.replace('"', '').replace('\\', '').replace('/', '_')[:100]
    encoded_filename = _url_quote(f"{safe_title}_{conv_id}.{format}")
    content_disposition = f"attachment; filename*=UTF-8''{encoded_filename}"

    if format == "txt":
        # 简单去 Markdown 标记
        import re as _re
        txt_content = _re.sub(r'[#*`~>\[\]\(\)]', '', md_content)
        return StreamingResponse(
            iter([txt_content]),
            media_type="text/plain; charset=utf-8",
            headers={"Content-Disposition": content_disposition}
        )
    return StreamingResponse(
        iter([md_content]),
        media_type="text/markdown; charset=utf-8",
        headers={"Content-Disposition": content_disposition}
    )


@app.get("/api/v1/chat/conversations/{conv_id}/messages")
async def chat_get_messages(conv_id: int, request: Request,
    limit: int = 200, offset: int = 0,
    user: dict = Depends(get_current_user)):
    """获取对话消息，支持分页；同时返回附件数据"""
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    msgs = await get_messages(conv_id, limit=limit, offset=offset)
    # 对回显的用户消息进行脱敏
    for m in msgs:
        if m.get("role") == "user":
            m["content"] = _mask_sensitive(m["content"])
    # 获取消息关联的附件
    from database import get_attachments_by_msg_ids
    msg_ids = [m["id"] for m in msgs]
    atts_by_msg = await get_attachments_by_msg_ids(conv_id, msg_ids) if msg_ids else {}
    # 附加到对应消息上
    for m in msgs:
        atts = atts_by_msg.get(m["id"], [])
        if atts:
            m["attachments"] = [{k: v for k, v in a.items() if k != "file_path"} for a in atts]
    return _success({"messages": msgs})


@app.put("/api/v1/chat/conversations/{conv_id}")
async def chat_rename_conversation(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """重命名对话"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("correct_json", request))
    title = str(body.get("title", "")).strip()
    if not title:
        return _error(400, _t("conversation_title_required", request))
    ok = await rename_conversation(conv_id, user["id"], title)
    if not ok:
        return _error(404, _t("conversation_not_found", request))
    return _success(None, "conversation_renamed", request)


@app.delete("/api/v1/chat/conversations/{conv_id}")
async def chat_delete_conversation(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """删除对话"""
    ok = await delete_conversation(conv_id, user["id"])
    if not ok:
        return _error(404, _t("conversation_not_found", request))
    return _success(None, "conversation_deleted", request)


@app.delete("/api/v1/chat/conversations/{conv_id}/messages/last")
async def chat_delete_last_message(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """删除对话最后一条 assistant 消息（用于重新生成）"""
    # 验证对话所有权
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    ok = await delete_last_assistant_message(conv_id)
    if not ok:
        return _success(None, "no_message_to_delete", request)
    return _success(None, "message_deleted", request)


@app.delete("/api/v1/chat/conversations/{conv_id}/messages")
async def chat_delete_message(conv_id: int, request: Request, msg_id: int = 0, user: dict = Depends(get_current_user)):
    """根据消息 ID 精确删除，避免索引竞态"""
    if not msg_id:
        return _error(400, _t("bad_request", request))
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    ok = await delete_message_by_id(conv_id, msg_id)
    if not ok:
        return _error(404, _t("message_not_found", request))
    return _success(None, "message_deleted", request)


# ---- 聊天附件上传 ----
ATTACHMENTS_DIR = os.path.join(os.path.dirname(__file__), "..", "uploads", "chat")
ATTACHMENT_MAX_SIZE = 10 * 1024 * 1024  # 10MB

@app.post("/api/v1/chat/conversations/{conv_id}/attachments")
async def chat_upload_attachment(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """上传聊天附件（图片/文档），返回 attachment_id 和预览信息"""
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))

    content_type = request.headers.get("content-type", "")
    if "multipart/form-data" not in content_type:
        return _error(400, "请使用 multipart/form-data 上传文件")

    form = await request.form()
    file = form.get("file")
    if not file:
        return _error(400, "未找到文件")

    filename = file.filename or "unknown"
    raw = await file.read()

    if len(raw) > ATTACHMENT_MAX_SIZE:
        return _error(400, f"文件过大（最大 10MB，当前 {len(raw)//1024//1024}MB）")

    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    is_image = ext in ('jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp')

    # 保存到日期子目录
    date_dir = datetime.now().strftime("%Y%m%d")
    save_dir = os.path.join(ATTACHMENTS_DIR, date_dir)
    os.makedirs(save_dir, exist_ok=True)
    safe_name = f"{uuid.uuid4().hex}_{filename}"
    save_path = os.path.join(save_dir, safe_name)
    with open(save_path, "wb") as f:
        f.write(raw)

    from database import add_attachment as db_add_attachment
    file_type = "image" if is_image else "document"
    preview = ""
    if is_image:
        preview = f"/uploads/chat/{date_dir}/{safe_name}"
    else:
        try:
            preview = _extract_doc_text(raw, ext, filename)[:300]
        except Exception:
            preview = f"[{filename}]"

    att_id = await db_add_attachment(conv_id, filename, file_type, save_path, preview)
    return _success({
        "id": att_id,
        "filename": filename,
        "file_type": file_type,
        "preview": preview,
        "url": f"/uploads/chat/{date_dir}/{safe_name}" if is_image else None,
    })


def _decode_text_file(raw: bytes) -> str:
    """解码文本文件：优先 UTF-8，若出现乱码则回退到 GB18030（兼容 GBK/GB2312）"""
    # 1. 去掉 UTF-8 BOM (如果存在)
    if raw.startswith(b'\xef\xbb\xbf'):
        raw = raw[3:]

    # 2. 尝试 UTF-8
    text = raw.decode('utf-8', errors='replace')
    if '\ufffd' not in text:
        return text

    # 3. UTF-8 解码出现替换字符 → 尝试 GB18030（GBK/GB2312 的超集）
    try:
        text_gb = raw.decode('gb18030')
        # 如果 GB18030 解出来的 CJK 字符更多，说明正确编码是 GB 系列
        utf8_cjk = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
        gb_cjk = sum(1 for c in text_gb if '\u4e00' <= c <= '\u9fff')
        if gb_cjk > utf8_cjk:
            return text_gb
        else:
            # 尝试 GBK（纯 GBK 更严格，有些字节序列在 GB18030 合法但在 GBK 不合法）
            text_gbk = raw.decode('gbk', errors='replace')
            gbk_cjk = sum(1 for c in text_gbk if '\u4e00' <= c <= '\u9fff')
            if gbk_cjk > utf8_cjk and '\ufffd' not in text_gbk:
                return text_gbk
    except (UnicodeDecodeError, UnicodeError):
        pass

    # 4. 都不行，返回 UTF-8 解码结果（至少有内容）
    return text



# OCR 暂不可用（PaddleOCR/PaddlePaddle 版本兼容问题，后续修复）
# _paddle_ocr = None
# _paddle_ocr_lock = None
#
# def _get_paddle_ocr():
#     ...


def _strip_code_fence(text: str) -> str:
    """剥离 LLM 输出的 markdown 代码围栏（```html ... ```）"""
    t = text.strip()
    # 去掉开头围栏（支持前导空白、大小写不敏感）
    m = re.match(r'\s*```\w*\s*\n?', t)
    if m:
        t = t[m.end():]
    # 去掉末尾围栏
    t = re.sub(r'\n?\s*```\s*$', '', t)
    return t.strip()


def _ocr_image(image_bytes: bytes) -> str:
    """OCR 暂不可用，始终返回空字符串"""
    return ""

def _extract_doc_text(raw: bytes, ext: str, filename: str) -> str:
    """从文档中提取全部文本（用于注入 LLM 上下文）
    支持: txt, md, csv, json, xml, yaml, log, docx, pdf, xlsx, xls, pptx"""
    if ext in ('txt', 'md', 'text', 'csv', 'json', 'xml', 'yaml', 'yml', 'log', 'ini', 'cfg', 'conf', 'toml', 'env',
               'py', 'js', 'ts', 'jsx', 'tsx', 'html', 'css', 'scss', 'less', 'sql', 'sh', 'bat', 'bash', 'zsh', 'ps1',
               'lua', 'cpp', 'cxx', 'cc', 'c', 'h', 'hpp', 'java', 'go', 'rs', 'rb', 'php', 'swift', 'kt', 'kts',
               'scala', 'pl', 'pm', 'r', 'm', 'mm', 'vue', 'svelte', 'tf', 'proto', 'cmake', 'make', 'dockerfile',
               'gradle', 'properties', 'gitignore', 'lock', 'patch', 'diff'):
        return _decode_text_file(raw)
    elif ext == 'docx':
        import io, docx
        doc = docx.Document(io.BytesIO(raw))
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(' | '.join(cells))
        return '\n'.join(parts)
    elif ext in ('xlsx', 'xlsm'):
        import io
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f'=== Sheet: {sheet_name} ===')
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(c.strip() for c in cells):
                        parts.append(' | '.join(cells))
                        row_count += 1
                        if row_count > 500:
                            parts.append('…(表格过大，已截断)')
                            break
            wb.close()
            return '\n'.join(parts)
        except ImportError:
            return f"[{filename}] (需要安装 openpyxl: pip install openpyxl)"
    elif ext == 'xls':
        import io
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=raw)
            parts = []
            for sheet_name in wb.sheet_names():
                ws = wb.sheet_by_name(sheet_name)
                parts.append(f'=== Sheet: {sheet_name} ===')
                for r in range(min(ws.nrows, 500)):
                    cells = [str(ws.cell_value(r, c)) if ws.cell_value(r, c) != '' else '' for c in range(ws.ncols)]
                    if any(c.strip() for c in cells):
                        parts.append(' | '.join(cells))
                if ws.nrows > 500:
                    parts.append('…(表格过大，已截断)')
            return '\n'.join(parts)
        except ImportError:
            return f"[{filename}] (需要安装 xlrd: pip install xlrd)"
    elif ext in ('pptx', 'ppt'):
        import io
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_texts.append(' | '.join(cells))
                if slide_texts:
                    parts.append(f'--- 第{i+1}页 ---')
                    parts.extend(slide_texts)
            return '\n'.join(parts)
        except ImportError:
            return f"[{filename}] (需要安装 python-pptx: pip install python-pptx)"
    elif ext == 'pdf':
        import fitz
        doc = fitz.open(stream=raw, filetype='pdf')
        parts = []
        for page in doc:
            t = page.get_text()
            if t and t.strip():
                parts.append(t)
        doc.close()
        text = '\n'.join(parts).strip()
        # 纯扫描件 PDF：文本为空时尝试 OCR
        if not text:
            doc = fitz.open(stream=raw, filetype='pdf')
            ocr_parts = []
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img_bytes = pix.tobytes("png")
                ocr_text = _ocr_image(img_bytes)
                if ocr_text:
                    ocr_parts.append(ocr_text)
            doc.close()
            return '\n'.join(ocr_parts) if ocr_parts else f"[{filename}] (无可提取文本)"
        return text
    elif ext in ('jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp', 'tiff', 'tif'):
        ocr_text = _ocr_image(raw)
        return ocr_text if ocr_text else f"[{filename}] (无可识别文字)"
    # 未知扩展名：尝试当作文本解码，失败则返回占位标记
    try:
        return _decode_text_file(raw)
    except Exception:
        return f"[{filename}]"


# ---- Fork 功能已禁用（前端未实现，后端端点预留） ----
# @app.post("/api/v1/chat/conversations/{conv_id}/messages")
# async def chat_add_message(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
#     ...


@app.post("/api/v1/chat/conversations/{conv_id}/messages/truncate")
async def chat_truncate_messages(conv_id: int, request: Request, user: dict = Depends(get_current_user)):
    """删除指定消息及其之后的所有消息，用于编辑消息后同步。
    优先使用 after_msg_id（基于消息 ID 精确截断）；兼容旧 after_index 参数。"""
    conv = await get_conversation(conv_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    after_msg_id = int(body.get("after_msg_id") or 0)
    strictly_after = bool(body.get("strictly_after", False))
    if after_msg_id > 0:
        deleted = await truncate_messages_from(conv_id, after_msg_id, inclusive=not strictly_after)
    else:
        # 兼容旧前端：使用 after_index（0-based 数组索引）
        after_index = int(body.get("after_index") or -1)
        if after_index < 0:
            return _error(400, _t("bad_request", request))
        # 将索引转换为消息 ID 后调用 truncate_messages_from
        msgs = await get_messages(conv_id)
        if after_index >= len(msgs):
            return _error(400, _t("bad_request", request))
        from_msg_id = msgs[after_index]["id"]
        deleted = await truncate_messages_from(conv_id, from_msg_id, inclusive=True)
    return _success({"deleted": deleted})


class ChatCompletionRequest(BaseModel):
    conversation_id: int
    content: str = Field(..., min_length=1, max_length=50000)
    regenerate: bool = False
    regenerate_msg_idx: int | None = None
    regenerate_msg_id: int | None = None
    attachment_ids: list[int] = []
    persona: str = "standard"


@app.post("/api/v1/chat/completions")
async def chat_completions(req: ChatCompletionRequest, request: Request, user: dict = Depends(get_current_user)):
    """流式对话：发送消息，接收流式回复"""
    await _check_rate_limit(user["id"], request.client.host, request)

    # 验证对话所有权
    conv = await get_conversation(req.conversation_id, user["id"])
    if not conv:
        return _error(404, _t("conversation_not_found", request))

    api_key, model, custom, deep_think = _get_headers(request)
    style = _get_style(request)
    chat_style_temp = _style_temperature(style)
    chat_style_extra = _style_instruction(style)

    # 关键词过滤 + 敏感信息脱敏
    _check_keywords(req.content, request)
    masked_content = _mask_sensitive(req.content)

    # regenerate: 前端已按 message ID 精确删除旧消息，后端不再重复截断。
    # 简单 regen（无 msg_idx）场景：删除最后一对 user+assistant 消息。
    if req.regenerate and req.regenerate_msg_idx is None:
        await delete_last_message_pair(req.conversation_id)

    # 验证 regenerate_msg_id 存在且属于该对话
    if req.regenerate and req.regenerate_msg_id:
        from database import _get_db
        db = await _get_db()
        c = await db.execute(
            "SELECT id FROM messages WHERE id = ? AND conversation_id = ?",
            (req.regenerate_msg_id, req.conversation_id),
        )
        if not await c.fetchone():
            # 消息已被删除（编辑后切到旧分支再生场景），降级为简单 regen
            logger.warning(f"regenerate_msg_id {req.regenerate_msg_id} not found, falling back to simple regen")
            req.regenerate_msg_id = None
            req.regenerate_msg_idx = None
            await delete_last_message_pair(req.conversation_id)

    # 加载历史消息
    history_messages = await get_messages(req.conversation_id)

    # 构建 messages 列表
    persona_id = req.persona or "standard"
    persona_system = get_chat_persona(persona_id)
    logger.info(f"Chat persona: {persona_id} (len={len(persona_system)})", extra={"request_id": getattr(request.state, "request_id", "-")})
    messages = [{"role": "system", "content": persona_system}]
    if custom:
        messages[0]["content"] = messages[0]["content"] + "\n\n【用户自定义指令】\n" + custom
    for msg in history_messages:
        # 剥离历史消息中的 KB 引用标记，避免 LLM 反复引用上一轮的文档
        import re as _re
        clean_content = _re.sub(r'【参考\s*\d+\s*·\s*来源：[^】]+】', '', msg["content"])
        messages.append({"role": msg["role"], "content": clean_content})

    # 附件处理：图片/文档文本前置注入（DeepSeek 不支持 multimodal，图片仅用文字标注）
    user_message_content = masked_content
    has_images = False
    if req.attachment_ids:
        attachments = await get_attachments(req.conversation_id, req.attachment_ids)
        image_notes = []
        doc_texts = []
        for att in attachments:
            if att["file_type"] == "image":
                has_images = True
                # 尝试 OCR 提取图片文字
                try:
                    with open(att["file_path"], "rb") as img_f:
                        img_raw = img_f.read()
                    ext = att["filename"].rsplit('.', 1)[-1].lower() if '.' in att["filename"] else ''
                    ocr_text = _extract_doc_text(img_raw, ext, att["filename"])
                    if ocr_text and not ocr_text.startswith('['):
                        logger.info(f"OCR success: {att['filename']} ({len(ocr_text)} chars)")
                        doc_texts.append(f"[图片文字识别：{att['filename']}]\n{ocr_text}")
                        continue
                    else:
                        logger.info(f"OCR no text: {att['filename']}")
                except Exception:
                    logger.warning(f"OCR extraction skipped")
                image_notes.append(
                    f"[用户上传了图片：{att['filename']}，OCR 未识别到文字。"
                    f"该图片可能是照片/截图/插画。请根据文件名和用户消息推断图片类型，"
                    f"引导用户描述图片内容，不要只说「看不到图片」]")
            else:
                # 从文件路径读取完整内容（而不是截断的 preview）
                try:
                    with open(att["file_path"], "rb") as doc_f:
                        doc_raw = doc_f.read()
                    ext = att["filename"].rsplit('.', 1)[-1].lower() if '.' in att["filename"] else ''
                    full_text = _extract_doc_text(doc_raw, ext, att["filename"])
                    # 限制单文件最多 20000 字避免撑爆上下文
                    if len(full_text) > 20000:
                        full_text = full_text[:20000] + "\n…(文档过长，后续内容已截断，建议分段上传)"
                    doc_texts.append(f"[文件：{att['filename']}]\n{full_text}")
                except Exception as e:
                    logger.warning(f"附件文本提取失败 file={att['file_path']} err={e}")
                    doc_texts.append(f"[文件：{att['filename']}]\n{att['preview']}")

        if image_notes or doc_texts:
            parts = []
            if doc_texts:
                parts.append("\n\n".join(doc_texts))
            if image_notes:
                parts.append("\n".join(image_notes))
            prefix = "\n\n".join(parts)
            user_message_content = prefix + "\n\n▔▔▔ 以上是文件内容，以下是用户消息 ▔▔▔\n" + masked_content
        messages.append({"role": "user", "content": user_message_content})
    else:
        messages.append({"role": "user", "content": masked_content})

    # 知识库检索注入（附件中有图片时跳过——消息是关于图片的，不是知识库查询）
    kb_chunks = []
    kb_col = request.headers.get("X-Kb-Collection", "")
    kb_cids = _parse_kb_collections(kb_col) if kb_col else None
    if kb_cids and not has_images and _need_kb_search(masked_content):
        try:
            keywords = await _extract_kb_keywords(masked_content, api_key, model)
            if keywords and _should_search_kb(keywords):
                from database import kb_search
                query = " OR ".join(keywords)
                chunks = await kb_search(user["id"], query, limit=KB_SEARCH_LIMIT, collection_ids=kb_cids)
                chunks = _filter_usable_chunks(chunks)
                if chunks and _kb_results_relevant(chunks, len(keywords)):
                    kb_chunks = [{"filename": c["filename"], "content": c["content"], "chunk_index": i}
                                 for i, c in enumerate(chunks)]
                    refs = "\n".join(f"【参考 {i+1} · 来源：{c['filename']}】\n{c['content'][:800]}"
                                    for i, c in enumerate(chunks))
                    messages[0]["content"] = messages[0]["content"] + "\n\n【知识库参考资料】\n" + refs + \
                        "\n\n以上资料根据关键词自动检索，仅供参考。请只关注用户本次发送的最新消息——只有当资料与本条消息直接相关时才引用。不要因为对话历史里聊过相关话题就引用资料，历史话题与当前问题可能已无关。引用格式：【参考 编号 · 来源：文件名】。如果不相关，直接忽略。" + \
                        "\n\n⚡ 注意：本轮对话中用户的问题主题可能已与之前不同。请以用户最新一条消息为准判断相关性，不要被历史话题干扰。"
        except Exception:
            logger.warning(f"KB search skipped in chat_completions")

    # 保存用户消息（脱敏后）
    # - 正常发送：总是保存
    # - regenerate 且给了 msg_idx（原位 regenerate）：用户消息仍在历史中，不保存
    # - regenerate 且无 msg_idx（简单 regen / 编辑重试）：delete_last_message_pair 已删除旧用户消息，需保存新消息
    user_msg_id = None
    save_user_msg = (not req.regenerate) or (req.regenerate and req.regenerate_msg_idx is None)
    if save_user_msg:
        user_msg_id = await add_message(req.conversation_id, "user", masked_content)

    # 将附件关联到该用户消息
    if req.attachment_ids and user_msg_id is not None:
        from database import link_attachments_to_message
        await link_attachments_to_message(req.conversation_id, user_msg_id, req.attachment_ids)

    async def _chat_stream():
        import json as _json
        full_response = ""
        assistant_msg_id = None
        _last_save_len = 0  # 上次保存时的长度
        client = None
        try:
            # 先创建占位 assistant 消息，拿到 msg_id 用于后续增量更新
            # regenerate 且提供了旧消息 ID → 原位更新，不新建
            if req.regenerate and req.regenerate_msg_id:
                assistant_msg_id = req.regenerate_msg_id
            else:
                assistant_msg_id = await add_message(req.conversation_id, "assistant", "")
            if not assistant_msg_id:
                yield f"data: {_json.dumps({'type':'error','message': _t('internal_error', request)}, ensure_ascii=False)}\n\n"
                return

            # 尽早回传 user_msg_id，让前端在 token 流到达前就能持久化版本历史
            if user_msg_id:
                yield f"data: {_json.dumps({'type':'user_msg','user_msg_id':user_msg_id}, ensure_ascii=False)}\n\n"

            key = api_key or ""
            model_name = model or ""

            # 深度思考模式：若未指定模型则用 reasoner；同时追加系统指令
            if deep_think:
                if not model:
                    model_name = "deepseek-reasoner"
                full_system = CHAT_SAFETY_GUARD + "\n\n" + messages[0]["content"]
                full_system += "\n\n【深度思考模式】请先进行逐步推理分析，再给出最终答案。"
            else:
                full_system = CHAT_SAFETY_GUARD + "\n\n" + messages[0]["content"]
            if chat_style_extra:
                full_system += chat_style_extra
            chat_messages = [{"role": "system", "content": full_system}] + messages[1:]

            client = AsyncOpenAI(
                api_key=key or os.getenv("DEEPSEEK_API_KEY", ""),
                base_url=ENV_BASE_URL,
                timeout=LLM_STREAM_TIMEOUT,
            )

            if is_mock_mode(key):
                mock_text = mock_response(messages[0]["content"], req.content)
                mock_text = _mask_sensitive(mock_text)  # 输出脱敏
                full_response = mock_text  # Mock 模式也需要累积完整响应用于保存
                for ch in mock_text:
                    yield f"data: {_json.dumps({'type':'token','content':ch}, ensure_ascii=False)}\n\n"
                    if ord(ch) % 3 == 0:
                        await asyncio.sleep(0.01)
                usage = {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0}
            else:
                response = await client.chat.completions.create(
                    model=model_name or os.getenv("DEEPSEEK_MODEL", "deepseek-chat"),
                    messages=chat_messages,
                    temperature=chat_style_temp,
                    max_tokens=20000,
                    stream=True,
                )
                prompt_tokens = 0
                completion_tokens = 0
                async for chunk in response:
                    delta = chunk.choices[0].delta if chunk.choices else None
                    if delta:
                        # 推理过程（仅深度思考模式开启时转发）
                        if deep_think and getattr(delta, 'reasoning_content', None):
                            thinking_text = delta.reasoning_content
                            yield f"data: {_json.dumps({'type':'thinking','content':thinking_text}, ensure_ascii=False)}\n\n"
                        if delta.content:
                            text = delta.content
                            full_response += text
                            # 每累积 ~200 字增量保存，防止刷新丢失
                            if assistant_msg_id and len(full_response) - _last_save_len >= 200:
                                try:
                                    await update_message_content(assistant_msg_id, full_response)
                                    _last_save_len = len(full_response)
                                except Exception as _save_err:
                                    logger.warning(f"Incremental save failed msg_id={assistant_msg_id}: {_save_err}")
                            yield f"data: {_json.dumps({'type':'token','content':text}, ensure_ascii=False)}\n\n"
                    if chunk.usage:
                        prompt_tokens = chunk.usage.prompt_tokens or 0
                        completion_tokens = chunk.usage.completion_tokens or 0
                usage = {
                    "prompt_tokens": prompt_tokens,
                    "completion_tokens": completion_tokens,
                    "total_tokens": prompt_tokens + completion_tokens,
                }

            # 保存 AI 回复（附带知识库引用片段 — 仅保存 LLM 实际引用的）
            if full_response.strip():
                full_response = _mask_sensitive(full_response)  # 输出脱敏
                # 只保存 LLM 实际引用的 KB 片段
                _kb_to_save = _filter_referenced_chunks(kb_chunks, full_response)
                kb_json = _json.dumps(_kb_to_save, ensure_ascii=False) if _kb_to_save else None
                # 最终写入完整内容（替换占位消息）
                if assistant_msg_id:
                    await update_message_content(assistant_msg_id, full_response, kb_json, final=True)
                elif full_response.strip():
                    await add_message(req.conversation_id, "assistant", full_response, kb_json)

            # 自动生成标题（首次对话完成时，用 LLM 摘要）
            if conv["title"] in ("新对话", ""):
                first_user_msg = next((m["content"] for m in history_messages if m["role"] == "user"), req.content)
                try:
                    new_title = await generate_conversation_title(
                        first_user_msg, _mask_sensitive(full_response),
                        api_key=key, model=model_name
                    )
                    if new_title and new_title not in ("新对话", ""):
                        await rename_conversation(req.conversation_id, user["id"], new_title)
                        yield f"data: {_json.dumps({'type':'title_update','title':new_title}, ensure_ascii=False)}\n\n"
                except Exception:
                    # 摘要失败则回退到简单截断（直接用第一条用户消息前 25 字）
                    fallback = (req.content or "新对话")[:25].strip()
                    if fallback and fallback not in ("新对话", ""):
                        await rename_conversation(req.conversation_id, user["id"], fallback)

            _kb_chunks_out = _filter_referenced_chunks(kb_chunks, full_response)
            yield f"data: {_json.dumps({'type':'done','msg_id':assistant_msg_id,'user_msg_id':user_msg_id,'usage':usage,'kb_chunks':_kb_chunks_out}, ensure_ascii=False)}\n\n"

        except (GeneratorExit, asyncio.CancelledError):
            # 客户端断开连接（刷新/关闭页面），保存当前已生成的部分内容
            if assistant_msg_id and full_response.strip():
                try:
                    await update_message_content(assistant_msg_id, _mask_sensitive(full_response))
                except Exception:
                    logger.warning("保存客户端断开时的部分内容失败")
            raise
        except Exception as e:
            # 发生错误时也保存已生成的部分
            if assistant_msg_id and full_response.strip():
                try:
                    await update_message_content(assistant_msg_id, _mask_sensitive(full_response))
                except Exception:
                    logger.warning("保存出错时的部分内容失败")
            logger.error(f"Chat stream error: {str(e)}", extra={"request_id": getattr(request.state, "request_id", "-")})
            yield f"data: {_json.dumps({'type':'error','message': _t('internal_error', request)}, ensure_ascii=False)}\n\n"
        finally:
            if client is not None:
                try:
                    await client.close()
                except Exception:
                    logger.warning("关闭 SSE 客户端连接失败")

    return StreamingResponse(
        _chat_stream(),
        media_type="text/event-stream",
        headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"}
    )


# ---- API Key 服务端加密存储 ----
@app.put("/api/v1/preferences/api-key")
async def save_api_key(request: Request, user: dict = Depends(get_current_user)):
    """将加密后的 API Key 存入用户偏好（加密由前端用 CryptoUtils 完成）"""
    try:
        body = await request.json()
    except Exception:
        return _error(400, _t("bad_request", request))
    encrypted_key = str(body.get("encrypted_key", "")).strip()
    if not encrypted_key:
        return _error(400, _t("bad_request", request))
    # 读取现有偏好，更新 api_key_enc 字段
    prefs_str = user.get("preferences", "{}")
    try:
        prefs = json.loads(prefs_str) if isinstance(prefs_str, str) else prefs_str
    except Exception:
        prefs = {}
    prefs["api_key_enc"] = encrypted_key
    from database import update_user_preferences
    ok = await update_user_preferences(user["id"], json.dumps(prefs, ensure_ascii=False))
    if not ok:
        return _error(500, _t("save_preferences_failed", request))
    return _success(None, "api_key_saved")


@app.get("/api/v1/preferences/api-key")
async def get_api_key(request: Request, user: dict = Depends(get_current_user)):
    """获取加密后的 API Key（前端用 CryptoUtils 解密）"""
    prefs_str = user.get("preferences", "{}")
    try:
        prefs = json.loads(prefs_str) if isinstance(prefs_str, str) else prefs_str
    except Exception:
        prefs = {}
    encrypted_key = prefs.get("api_key_enc", "")
    return _success({"encrypted_key": encrypted_key})


# ---- 静态文件 ----
_uploads_dir = os.path.join(os.path.dirname(__file__), "..", "uploads")
os.makedirs(_uploads_dir, exist_ok=True)

# 受保护的 uploads 文件服务（需登录，支持 Header 和 Cookie 两种鉴权方式，
# Cookie 鉴权使 <img> 标签等无法自定义 Header 的场景也能正常工作）
@app.get("/uploads/{file_path:path}")
async def serve_upload(file_path: str, request: Request):
    """鉴权文件下载：所有 /uploads/ 下的文件均需登录后才能访问"""
    import mimetypes
    from fastapi.responses import FileResponse
    # 鉴权：优先 Authorization Header，回退到 Cookie
    token = None
    auth_header = request.headers.get("Authorization", "")
    if auth_header.startswith("Bearer "):
        token = auth_header[7:]
    if not token:
        token = request.cookies.get("grindpal_token", "")
    if not token:
        return _error(401, _t("please_login", request))
    try:
        from auth import decode_token as _decode
        payload = _decode(token)
    except Exception:
        return _error(401, _t("login_expired", request))
    safe_path = os.path.normpath(os.path.join(_uploads_dir, file_path))
    # 路径穿越防护
    if not safe_path.startswith(os.path.normpath(_uploads_dir) + os.sep) and safe_path != os.path.normpath(_uploads_dir):
        return _error(403, _t("forbidden", request))
    if not os.path.isfile(safe_path):
        return _error(404, _t("file_not_found", request))
    media_type, _ = mimetypes.guess_type(safe_path)
    return FileResponse(safe_path, media_type=media_type or "application/octet-stream")
_frontend_dir = os.path.join(os.path.dirname(__file__), "..", "frontend")
if os.path.isdir(_frontend_dir):
    app.mount("/", StaticFiles(directory=_frontend_dir, html=True), name="frontend")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app",
        host=os.getenv("HOST", "0.0.0.0"),
        port=int(os.getenv("PORT", "8000")),
        reload=False)
